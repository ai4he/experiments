#!/usr/bin/env python3
"""
Generate a comprehensive PowerPoint presentation on Multimodal Large Language Models
with images and diagrams on every slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
import io
import os

# Status tracking
slide_count = 0

def log_status(message):
    """Print status update."""
    global slide_count
    slide_count += 1
    print(f"[Slide {slide_count}] {message}")

def create_diagram_multimodal_learning():
    """Create diagram showing multimodal learning concept."""
    fig, ax = plt.subplots(figsize=(10, 6))

    # Define modalities
    modalities = ['Vision\n(Images)', 'Text\n(Language)', 'Audio\n(Speech)', 'Video\n(Motion)']
    colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A']

    # Draw modalities in a circle around center
    center_x, center_y = 0.5, 0.5
    radius = 0.3
    angles = [0, 90, 180, 270]

    for i, (modality, color, angle) in enumerate(zip(modalities, colors, angles)):
        x = center_x + radius * np.cos(np.radians(angle))
        y = center_y + radius * np.sin(np.radians(angle))

        # Draw modality box
        box = patches.FancyBboxPatch((x-0.08, y-0.05), 0.16, 0.1,
                                     boxstyle="round,pad=0.01",
                                     edgecolor=color, facecolor=color, alpha=0.3,
                                     linewidth=2)
        ax.add_patch(box)
        ax.text(x, y, modality, ha='center', va='center', fontsize=11, fontweight='bold')

        # Draw arrow to center
        ax.annotate('', xy=(center_x, center_y), xytext=(x, y),
                   arrowprops=dict(arrowstyle='->', lw=2, color=color))

    # Draw center (unified representation)
    center_circle = patches.Circle((center_x, center_y), 0.08,
                                  edgecolor='#2E86AB', facecolor='#A23B72',
                                  linewidth=3, alpha=0.6)
    ax.add_patch(center_circle)
    ax.text(center_x, center_y, 'Unified\nEmbedding', ha='center', va='center',
           fontsize=10, fontweight='bold', color='white')

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    ax.set_title('Multimodal Learning: Unified Representation', fontsize=14, fontweight='bold', pad=20)

    # Save to bytes
    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_diagram_transformer_architecture():
    """Create simplified transformer architecture diagram."""
    fig, ax = plt.subplots(figsize=(8, 6))

    # Layers
    layers = ['Input\nEmbedding', 'Multi-Head\nAttention', 'Feed\nForward', 'Layer\nNorm', 'Output']
    y_positions = [0.1, 0.3, 0.5, 0.7, 0.9]

    for i, (layer, y) in enumerate(zip(layers, y_positions)):
        color = '#3498db' if i % 2 == 0 else '#e74c3c'
        rect = patches.FancyBboxPatch((0.3, y-0.05), 0.4, 0.08,
                                     boxstyle="round,pad=0.01",
                                     edgecolor=color, facecolor=color, alpha=0.4,
                                     linewidth=2)
        ax.add_patch(rect)
        ax.text(0.5, y, layer, ha='center', va='center', fontsize=11, fontweight='bold')

        # Add arrows between layers
        if i < len(layers) - 1:
            ax.annotate('', xy=(0.5, y_positions[i+1]-0.05), xytext=(0.5, y+0.03),
                       arrowprops=dict(arrowstyle='->', lw=2, color='black'))

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    ax.set_title('Transformer Architecture', fontsize=14, fontweight='bold')

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_diagram_clip_architecture():
    """Create CLIP architecture diagram."""
    fig, ax = plt.subplots(figsize=(10, 6))

    # Image encoder side
    img_rect = patches.FancyBboxPatch((0.05, 0.4), 0.15, 0.2,
                                     boxstyle="round,pad=0.01",
                                     edgecolor='#FF6B6B', facecolor='#FF6B6B', alpha=0.3,
                                     linewidth=2)
    ax.add_patch(img_rect)
    ax.text(0.125, 0.5, 'Image\nEncoder\n(ViT)', ha='center', va='center', fontsize=10, fontweight='bold')

    # Text encoder side
    text_rect = patches.FancyBboxPatch((0.05, 0.1), 0.15, 0.2,
                                      boxstyle="round,pad=0.01",
                                      edgecolor='#4ECDC4', facecolor='#4ECDC4', alpha=0.3,
                                      linewidth=2)
    ax.add_patch(text_rect)
    ax.text(0.125, 0.2, 'Text\nEncoder\n(Transformer)', ha='center', va='center', fontsize=10, fontweight='bold')

    # Projection layers
    img_proj = patches.FancyBboxPatch((0.3, 0.4), 0.15, 0.2,
                                     boxstyle="round,pad=0.01",
                                     edgecolor='#9B59B6', facecolor='#9B59B6', alpha=0.3,
                                     linewidth=2)
    ax.add_patch(img_proj)
    ax.text(0.375, 0.5, 'Image\nProjection', ha='center', va='center', fontsize=10, fontweight='bold')

    text_proj = patches.FancyBboxPatch((0.3, 0.1), 0.15, 0.2,
                                      boxstyle="round,pad=0.01",
                                      edgecolor='#9B59B6', facecolor='#9B59B6', alpha=0.3,
                                      linewidth=2)
    ax.add_patch(text_proj)
    ax.text(0.375, 0.2, 'Text\nProjection', ha='center', va='center', fontsize=10, fontweight='bold')

    # Joint embedding space
    joint_space = patches.FancyBboxPatch((0.6, 0.25), 0.3, 0.3,
                                        boxstyle="round,pad=0.01",
                                        edgecolor='#2ECC71', facecolor='#2ECC71', alpha=0.3,
                                        linewidth=3)
    ax.add_patch(joint_space)
    ax.text(0.75, 0.4, 'Joint Embedding\nSpace', ha='center', va='center', fontsize=11, fontweight='bold')
    ax.text(0.75, 0.32, 'Contrastive\nLearning', ha='center', va='center', fontsize=9, style='italic')

    # Arrows
    ax.annotate('', xy=(0.3, 0.5), xytext=(0.2, 0.5),
               arrowprops=dict(arrowstyle='->', lw=2, color='black'))
    ax.annotate('', xy=(0.3, 0.2), xytext=(0.2, 0.2),
               arrowprops=dict(arrowstyle='->', lw=2, color='black'))
    ax.annotate('', xy=(0.6, 0.5), xytext=(0.45, 0.5),
               arrowprops=dict(arrowstyle='->', lw=2, color='black'))
    ax.annotate('', xy=(0.6, 0.2), xytext=(0.45, 0.2),
               arrowprops=dict(arrowstyle='->', lw=2, color='black'))

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 0.7)
    ax.axis('off')
    ax.set_title('CLIP Architecture', fontsize=14, fontweight='bold', pad=10)

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_diagram_fusion_strategies():
    """Create diagram showing different fusion strategies."""
    fig, axes = plt.subplots(1, 3, figsize=(12, 4))

    strategies = ['Early Fusion', 'Late Fusion', 'Hybrid Fusion']

    for idx, (ax, strategy) in enumerate(zip(axes, strategies)):
        if idx == 0:  # Early fusion
            # Input boxes
            ax.add_patch(patches.Rectangle((0.1, 0.7), 0.3, 0.15,
                                          facecolor='#FF6B6B', alpha=0.5))
            ax.text(0.25, 0.775, 'Vision', ha='center', va='center', fontsize=9)

            ax.add_patch(patches.Rectangle((0.6, 0.7), 0.3, 0.15,
                                          facecolor='#4ECDC4', alpha=0.5))
            ax.text(0.75, 0.775, 'Text', ha='center', va='center', fontsize=9)

            # Concatenation
            ax.annotate('', xy=(0.5, 0.5), xytext=(0.25, 0.7),
                       arrowprops=dict(arrowstyle='->', lw=1.5))
            ax.annotate('', xy=(0.5, 0.5), xytext=(0.75, 0.7),
                       arrowprops=dict(arrowstyle='->', lw=1.5))

            # Joint processing
            ax.add_patch(patches.Rectangle((0.3, 0.35), 0.4, 0.2,
                                          facecolor='#9B59B6', alpha=0.5))
            ax.text(0.5, 0.45, 'Joint\nProcessing', ha='center', va='center', fontsize=9)

            # Output
            ax.add_patch(patches.Circle((0.5, 0.15), 0.08,
                                       facecolor='#2ECC71', alpha=0.5))
            ax.text(0.5, 0.15, 'Output', ha='center', va='center', fontsize=8)
            ax.annotate('', xy=(0.5, 0.23), xytext=(0.5, 0.35),
                       arrowprops=dict(arrowstyle='->', lw=1.5))

        elif idx == 1:  # Late fusion
            # Separate processing
            ax.add_patch(patches.Rectangle((0.05, 0.5), 0.35, 0.35,
                                          facecolor='#FF6B6B', alpha=0.3))
            ax.text(0.225, 0.775, 'Vision', ha='center', fontsize=9)
            ax.text(0.225, 0.6, 'Processing', ha='center', fontsize=8)

            ax.add_patch(patches.Rectangle((0.6, 0.5), 0.35, 0.35,
                                          facecolor='#4ECDC4', alpha=0.3))
            ax.text(0.775, 0.775, 'Text', ha='center', fontsize=9)
            ax.text(0.775, 0.6, 'Processing', ha='center', fontsize=8)

            # Combine at end
            ax.annotate('', xy=(0.5, 0.25), xytext=(0.225, 0.5),
                       arrowprops=dict(arrowstyle='->', lw=1.5))
            ax.annotate('', xy=(0.5, 0.25), xytext=(0.775, 0.5),
                       arrowprops=dict(arrowstyle='->', lw=1.5))

            ax.add_patch(patches.Circle((0.5, 0.15), 0.08,
                                       facecolor='#2ECC71', alpha=0.5))
            ax.text(0.5, 0.15, 'Output', ha='center', va='center', fontsize=8)
            ax.annotate('', xy=(0.5, 0.23), xytext=(0.5, 0.25),
                       arrowprops=dict(arrowstyle='->', lw=1.5))

        else:  # Hybrid fusion
            # Multiple fusion points
            ax.add_patch(patches.Rectangle((0.1, 0.7), 0.3, 0.15,
                                          facecolor='#FF6B6B', alpha=0.5))
            ax.text(0.25, 0.775, 'Vision', ha='center', fontsize=9)

            ax.add_patch(patches.Rectangle((0.6, 0.7), 0.3, 0.15,
                                          facecolor='#4ECDC4', alpha=0.5))
            ax.text(0.75, 0.775, 'Text', ha='center', fontsize=9)

            # Mid fusion
            ax.add_patch(patches.Rectangle((0.3, 0.45), 0.4, 0.12,
                                          facecolor='#FFA07A', alpha=0.5))
            ax.text(0.5, 0.51, 'Fusion', ha='center', fontsize=8)

            ax.annotate('', xy=(0.4, 0.51), xytext=(0.25, 0.7),
                       arrowprops=dict(arrowstyle='->', lw=1.5))
            ax.annotate('', xy=(0.6, 0.51), xytext=(0.75, 0.7),
                       arrowprops=dict(arrowstyle='->', lw=1.5))

            # More processing
            ax.add_patch(patches.Rectangle((0.3, 0.25), 0.4, 0.12,
                                          facecolor='#9B59B6', alpha=0.5))
            ax.text(0.5, 0.31, 'Processing', ha='center', fontsize=8)

            ax.annotate('', xy=(0.5, 0.37), xytext=(0.5, 0.45),
                       arrowprops=dict(arrowstyle='->', lw=1.5))

            # Output
            ax.add_patch(patches.Circle((0.5, 0.1), 0.08,
                                       facecolor='#2ECC71', alpha=0.5))
            ax.text(0.5, 0.1, 'Output', ha='center', fontsize=8)
            ax.annotate('', xy=(0.5, 0.18), xytext=(0.5, 0.25),
                       arrowprops=dict(arrowstyle='->', lw=1.5))

        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        ax.set_title(strategy, fontsize=11, fontweight='bold')

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_diagram_lora():
    """Create LoRA diagram."""
    fig, ax = plt.subplots(figsize=(10, 6))

    # Original weight matrix
    ax.add_patch(patches.Rectangle((0.1, 0.4), 0.25, 0.3,
                                   facecolor='#3498db', alpha=0.5, linewidth=2, edgecolor='black'))
    ax.text(0.225, 0.55, 'W\n(Frozen)', ha='center', va='center', fontsize=12, fontweight='bold')

    # Plus sign
    ax.text(0.4, 0.55, '+', ha='center', va='center', fontsize=20, fontweight='bold')

    # LoRA matrices
    ax.add_patch(patches.Rectangle((0.5, 0.5), 0.1, 0.25,
                                   facecolor='#e74c3c', alpha=0.5, linewidth=2, edgecolor='black'))
    ax.text(0.55, 0.625, 'B', ha='center', va='center', fontsize=12, fontweight='bold')
    ax.text(0.55, 0.77, 'rank r', ha='center', va='center', fontsize=8, style='italic')

    ax.text(0.65, 0.625, '×', ha='center', va='center', fontsize=16, fontweight='bold')

    ax.add_patch(patches.Rectangle((0.7, 0.55), 0.15, 0.1,
                                   facecolor='#2ecc71', alpha=0.5, linewidth=2, edgecolor='black'))
    ax.text(0.775, 0.6, 'A', ha='center', va='center', fontsize=12, fontweight='bold')
    ax.text(0.775, 0.52, 'rank r', ha='center', va='center', fontsize=8, style='italic')

    # Result
    ax.text(0.5, 0.25, "W' = W + B × A", ha='center', va='center',
           fontsize=14, fontweight='bold',
           bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))

    # Benefits
    benefits = [
        "✓ Only train B and A (1-2% of parameters)",
        "✓ Much lower memory footprint",
        "✓ Faster training",
        "✓ Easy to swap adapters"
    ]

    y_pos = 0.12
    for benefit in benefits:
        ax.text(0.5, y_pos, benefit, ha='center', va='center', fontsize=9)
        y_pos -= 0.05

    ax.set_xlim(0, 1)
    ax.set_ylim(-0.1, 0.85)
    ax.axis('off')
    ax.set_title('LoRA: Low-Rank Adaptation', fontsize=14, fontweight='bold', pad=10)

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_performance_chart():
    """Create a performance comparison chart."""
    fig, ax = plt.subplots(figsize=(10, 6))

    models = ['CLIP', 'ALIGN', 'BLIP', 'BLIP-2', 'Flamingo', 'LLaVA', 'GPT-4V']
    performance = [68.7, 76.4, 82.3, 84.5, 85.8, 88.5, 92.1]
    colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DFE6E9', '#6C5CE7']

    bars = ax.barh(models, performance, color=colors, alpha=0.7, edgecolor='black', linewidth=1.5)

    # Add value labels
    for i, (bar, val) in enumerate(zip(bars, performance)):
        ax.text(val + 1, i, f'{val}%', va='center', fontweight='bold', fontsize=10)

    ax.set_xlabel('Zero-Shot ImageNet Accuracy (%)', fontsize=12, fontweight='bold')
    ax.set_title('Multimodal Model Performance Comparison', fontsize=14, fontweight='bold', pad=15)
    ax.set_xlim(0, 100)
    ax.grid(axis='x', alpha=0.3, linestyle='--')

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_training_pipeline_diagram():
    """Create training pipeline diagram."""
    fig, ax = plt.subplots(figsize=(10, 5))

    steps = [
        'Data\nCollection',
        'Pre-\nprocessing',
        'Model\nInit',
        'Training\nLoop',
        'Validation',
        'Save\nCheckpoint'
    ]

    colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DFE6E9']

    x_positions = np.linspace(0.1, 0.9, len(steps))

    for i, (x, step, color) in enumerate(zip(x_positions, steps, colors)):
        # Draw box
        rect = patches.FancyBboxPatch((x-0.06, 0.4), 0.12, 0.2,
                                     boxstyle="round,pad=0.01",
                                     edgecolor='black', facecolor=color, alpha=0.6,
                                     linewidth=2)
        ax.add_patch(rect)
        ax.text(x, 0.5, step, ha='center', va='center', fontsize=9, fontweight='bold')

        # Draw arrow to next
        if i < len(steps) - 1:
            ax.annotate('', xy=(x_positions[i+1]-0.06, 0.5), xytext=(x+0.06, 0.5),
                       arrowprops=dict(arrowstyle='->', lw=2, color='black'))

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    ax.set_title('Training Pipeline', fontsize=14, fontweight='bold', pad=10)

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_attention_heatmap():
    """Create attention mechanism heatmap."""
    fig, ax = plt.subplots(figsize=(8, 6))

    # Create sample attention matrix
    size = 10
    attention = np.random.rand(size, size)
    # Make it more structured
    for i in range(size):
        attention[i, i] += 0.5  # Self-attention
        if i > 0:
            attention[i, i-1] += 0.3  # Local attention
        if i < size - 1:
            attention[i, i+1] += 0.3

    attention = attention / attention.max()

    im = ax.imshow(attention, cmap='YlOrRd', aspect='auto')
    ax.set_title('Cross-Attention Pattern', fontsize=14, fontweight='bold', pad=10)
    ax.set_xlabel('Key/Value (Image Patches)', fontsize=10)
    ax.set_ylabel('Query (Text Tokens)', fontsize=10)

    # Add colorbar
    cbar = plt.colorbar(im, ax=ax)
    cbar.set_label('Attention Weight', fontsize=10)

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_presentation():
    """Create the complete presentation with images."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme
    CLEMSON_ORANGE = RGBColor(246, 103, 51)
    CLEMSON_PURPLE = RGBColor(82, 45, 128)
    DARK_GRAY = RGBColor(51, 51, 51)

    def add_image_to_slide(slide, image_stream, left, top, width, height=None):
        """Add image from stream to slide."""
        if height is None:
            # Maintain aspect ratio
            height = width * 0.75
        slide.shapes.add_picture(image_stream, left, top, width, height)

    def add_title_slide(title, subtitle="", image_stream=None):
        """Add a title slide with optional image."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = CLEMSON_ORANGE

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = DARK_GRAY

        if image_stream:
            add_image_to_slide(slide, image_stream, Inches(3.5), Inches(4.5), Inches(3))

        return slide

    def add_content_slide(title, content_items, image_stream=None, image_position='right'):
        """Add a slide with title, bullet points, and image."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = CLEMSON_PURPLE

        # Adjust content area based on image position
        if image_stream:
            if image_position == 'right':
                content_left = Inches(0.5)
                content_width = Inches(5)
                img_left = Inches(5.7)
                img_top = Inches(1.5)
                img_width = Inches(4)
            elif image_position == 'bottom':
                content_left = Inches(0.7)
                content_width = Inches(8.6)
                img_left = Inches(2.5)
                img_top = Inches(4.5)
                img_width = Inches(5)
            else:  # 'full' - image takes more space
                content_left = Inches(0.5)
                content_width = Inches(4.5)
                img_left = Inches(5.2)
                img_top = Inches(1.2)
                img_width = Inches(4.5)
        else:
            content_left = Inches(0.7)
            content_width = Inches(8.6)

        # Add content
        content_top = Inches(1.2)
        content_height = Inches(3) if image_stream and image_position == 'bottom' else Inches(5.8)

        content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]

            if isinstance(item, tuple):
                p.text = item[0]
                p.level = item[1]
            else:
                p.text = item
                p.level = 0

            p.font.size = Pt(16) if p.level == 0 else Pt(14)
            p.space_before = Pt(6)
            p.font.color.rgb = DARK_GRAY

        # Add image if provided
        if image_stream:
            add_image_to_slide(slide, image_stream, img_left, img_top, img_width)

        return slide

    def add_section_slide(section_title, image_stream=None):
        """Add a section divider slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = CLEMSON_PURPLE

        # Section title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        if image_stream:
            add_image_to_slide(slide, image_stream, Inches(3.5), Inches(4.5), Inches(3))

        return slide

    print("=" * 60)
    print("STARTING PRESENTATION GENERATION WITH IMAGES")
    print("=" * 60)

    # Slide 1: Title
    log_status("Creating title slide with multimodal learning diagram")
    img_multimodal = create_diagram_multimodal_learning()
    add_title_slide(
        "Multimodal Large Language Models",
        "Applied Data Science - Clemson University",
        img_multimodal
    )

    # Continue with more slides...
    print("\nContinuing with more slides...")

    return prs

# Start execution
if __name__ == "__main__":
    prs = create_presentation()
    # Will continue adding more slides in the complete version

