#!/usr/bin/env python3
"""Add slide 56"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-55.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

print("[Slide 56/78] Adding: Model Deployment Strategies")

img = create_simple_diagram("Deployment", ["Local\nInference", "API\nService", "Edge\nDevices", "Cloud"])
add_comprehensive_slide("Model Deployment Strategies", [
    "Deployment Overview and Considerations: Deploying multimodal models in production requires balancing multiple objectives: inference latency (response time), throughput (requests per second), computational cost (hardware and energy), model size (memory footprint), and accuracy. Different deployment scenarios have different constraints: (1) Interactive applications (chatbots, image search) require low latency (<100ms) and high availability. (2) Batch processing (content moderation, dataset annotation) prioritizes throughput over latency. (3) Edge deployment (mobile apps, robotics) has strict memory and compute constraints. (4) Research and development prioritizes flexibility and experimentation. Key decisions: where to deploy (cloud, on-premise, edge), what hardware (CPU, GPU, TPU, specialized accelerators), what optimizations (quantization, pruning, distillation), and what serving infrastructure (REST API, gRPC, model server frameworks).",

    "Cloud Deployment with Model Serving Frameworks: Cloud deployment provides scalability and flexibility. Popular frameworks: (1) TorchServe - PyTorch's native serving framework supporting batching, multi-model serving, metrics, and A/B testing. Package model as .mar file, deploy with configuration specifying batch size, timeout, workers. Handles HTTP/gRPC requests, load balancing, auto-scaling. (2) TensorFlow Serving - similar for TensorFlow models. (3) Triton Inference Server (NVIDIA) - framework-agnostic, optimized for GPU inference, supports PyTorch, TensorFlow, ONNX. Features: dynamic batching (accumulating requests to maximize GPU utilization), model ensembles, concurrent model execution. (4) Ray Serve - Python-first serving framework with advanced features like multi-model composition, streaming, custom business logic. Deployment workflow: Export model to standard format (TorchScript, ONNX), containerize with Docker, deploy to Kubernetes cluster with horizontal auto-scaling, monitor with Prometheus/Grafana.",

    "Edge Deployment and Mobile Optimization: Deploying on edge devices (smartphones, IoT, robots) requires aggressive optimization. Challenges: limited memory (2-8GB), CPU-only or mobile GPU, battery constraints, intermittent connectivity. Optimization techniques: (1) Quantization - convert FP32 to INT8 or even INT4 using post-training quantization or quantization-aware training. PyTorch Mobile and TensorFlow Lite provide quantization tools. (2) Pruning - remove unimportant weights/neurons, reducing model size and computation. Structured pruning removes entire channels/layers. (3) Knowledge distillation - train smaller student model to mimic larger teacher. For example, distilling CLIP ViT-L to ViT-B or MobileNet backbone. (4) Neural architecture search - finding efficient architectures for target hardware (MobileNets, EfficientNets designed for mobile). (5) Model partitioning - split computation between edge and cloud, processing sensitive data locally while offloading heavy computation. Example: Run CLIP's vision encoder on device, text encoder in cloud for privacy-preserving image search.",

    "Inference Optimization and Best Practices: Optimizations applicable across deployment scenarios: (1) Batching - processing multiple inputs together amortizes overhead and maximizes hardware utilization. Dynamic batching accumulates requests up to timeout or max batch size. Critical for GPU throughput. (2) Model compilation - JIT compilation (TorchScript torch.jit.trace/script) optimizes computation graph. ONNX Runtime provides further optimizations. TensorRT (NVIDIA) aggressively optimizes for specific GPU models. (3) Mixed precision inference - FP16 reduces memory and increases throughput with minimal accuracy loss. (4) Caching - cache embeddings for frequently accessed items (e.g., cache CLIP image embeddings for product catalog). (5) Multi-stage pipelines - for complex workflows, use smaller fast models for filtering before expensive models. (6) Hardware selection - GPUs for throughput-critical applications, CPUs with AVX-512/VNNI for cost-sensitive deployment, TPUs for Google Cloud. Monitoring: Track latency percentiles (p50, p95, p99), throughput, error rates, resource utilization. Set up alerts for degradation. A/B testing for model updates. Gradual rollout to catch issues before full deployment."
], img)

output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-56.pptx'
prs.save(output_file)
print(f"✓ Slide 56 completed: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
