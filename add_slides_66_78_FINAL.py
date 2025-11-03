#!/usr/bin/env python3
"""Add final slides 66-78 to complete the presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
CLEMSON_ORANGE = RGBColor(246, 103, 51)
DARK_GRAY = RGBColor(51, 51, 51)

prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-65.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = CLEMSON_PURPLE
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    return slide

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_ORANGE
    if subtitle:
        stb = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
        stf = stb.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sp.font.size = Pt(22)
        sp.font.color.rgb = DARK_GRAY
    return slide

print("[Slide 66/78] Section: Conclusion")
add_section_slide("Part VII: Conclusion and Next Steps")

print("[Slide 67/78] Key Takeaways")
img = create_simple_diagram("Takeaways", ["Theory", "Practice", "Ethics", "Future"])
add_comprehensive_slide("Key Takeaways", [
    "Multimodal learning integrates information across modalities (vision, language, audio) to build richer AI systems than single-modality approaches. Key theoretical principles: contrastive learning aligns embeddings across modalities via InfoNCE loss, attention mechanisms enable flexible cross-modal information flow, and transfer learning leverages pre-trained models to reduce data and compute requirements. These foundations enable zero-shot capabilities, efficient adaptation, and emergent cross-modal reasoning absent in modality-specific models.",

    "The architectural landscape offers diverse options with different trade-offs. Dual encoders (CLIP) provide computational efficiency and modularity. Cross-encoders (BLIP) enable rich interaction but at higher cost. Frozen-backbone models (BLIP-2, Flamingo, LLaVA) achieve strong performance with parameter-efficient training. No single architecture dominates - choose based on task requirements, data availability, and computational constraints. Understanding these trade-offs enables informed architectural decisions in research and applications.",

    "Practical deployment requires balancing multiple objectives: accuracy, latency, throughput, cost, fairness, and safety. Techniques like quantization, LoRA fine-tuning, and optimized serving infrastructure enable deploying large models in resource-constrained environments. Comprehensive testing, monitoring, and documentation ensure reliable production systems. Ethical considerations - bias mitigation, privacy protection, transparency - are not afterthoughts but integral to responsible AI development and deployment.",

    "The field is rapidly evolving. Emerging directions include 3D vision, long-form video understanding, embodied AI for robotics, and domain-specific models for healthcare, science, and industry. Open challenges remain in reasoning, hallucination mitigation, computational efficiency, and safety. Your generation of researchers and engineers will shape the future of multimodal AI. The foundations covered in this lecture provide a starting point for contributing to this exciting field."
], img)

print("[Slide 68/78] Recommended Reading")
img = create_simple_diagram("Reading", ["Papers", "Books", "Blogs", "Courses"])
add_comprehensive_slide("Recommended Reading and Learning Resources", [
    "Foundational Papers (Must Read): (1) 'Attention Is All You Need' - Vaswani et al., NeurIPS 2017. Introduces Transformer architecture underlying all modern multimodal models. (2) 'CLIP: Learning Transferable Visual Models From Natural Language Supervision' - Radford et al., ICML 2021. Seminal work on contrastive vision-language learning and zero-shot capabilities. (3) 'BLIP-2: Bootstrapping Language-Image Pre-training' - Li et al., ICML 2023. Efficient architecture connecting frozen models. (4) 'Flamingo: a Visual Language Model for Few-Shot Learning' - Alayrac et al., NeurIPS 2022. In-context learning for multimodal tasks. (5) 'LoRA: Low-Rank Adaptation' - Hu et al., ICLR 2022. Parameter-efficient fine-tuning enabling adaptation on modest hardware.",

    "Books and Comprehensive Resources: (1) 'Deep Learning' by Goodfellow, Bengio, Courville - foundational textbook covering neural networks, optimization, generative models. Free online at deeplearningbook.org. (2) 'Dive into Deep Learning' - interactive textbook with code examples in PyTorch/TensorFlow/JAX. d2l.ai. (3) 'Speech and Language Processing' by Jurafsky and Martin - comprehensive NLP textbook. Draft chapters free online. (4) 'Computer Vision: Algorithms and Applications' by Szeliski - thorough computer vision reference. (5) Stanford CS231n (Computer Vision), CS224n (NLP), CS229 (Machine Learning) lecture notes and videos freely available.",

    "Blogs and Online Content: (1) Distill.pub - high-quality visualizations and explanations of ML concepts. (2) The Gradient - ML research magazine with accessible articles. (3) HuggingFace Blog - tutorials, model releases, research discussions. (4) Jay Alammar's Blog - excellent visual explanations of Transformers, BERT, GPT. (5) Lil'Log by Lilian Weng (OpenAI) - in-depth technical posts on recent research. (6) Google AI Blog and Meta AI Blog for industry research updates. (7) YouTube: Yannic Kilcher (paper explanations), Two Minute Papers (research highlights), StatQuest (statistics foundations). (8) Podcasts: The TWIML AI Podcast, Gradient Dissent, Practical AI.",

    "Hands-On Learning: (1) Fast.ai courses - practical deep learning with top-down approach. (2) HuggingFace Course - transformers, datasets, and deployment. (3) Coursera: Andrew Ng's ML Specialization, Deep Learning Specialization. (4) Kaggle Learn - bite-sized tutorials on ML topics plus competitions for practice. (5) Papers With Code - find implementations of papers, compare benchmarks. (6) Build projects - nothing substitutes hands-on implementation. Start with reproducing paper results, then extend to new ideas. (7) Contribute to open source - HuggingFace, PyTorch, TensorFlow welcome contributions. (8) Write blog posts explaining concepts - teaching solidifies understanding. (9) Participate in reading groups or study groups - discussing papers deepens comprehension."
], img)

print("[Slide 69/78] Tools and Libraries Reference")
img = create_simple_diagram("Tools", ["PyTorch", "HuggingFace", "Datasets", "MLOps"])
add_comprehensive_slide("Essential Tools and Libraries", [
    "Core Deep Learning Frameworks: PyTorch (pytorch.org) - most popular for research, dynamic computation graphs, extensive ecosystem. TensorFlow (tensorflow.org) - strong for production deployment, TF Lite for mobile, TF Serving for serving. JAX (jax.readthedocs.io) - composable transformations (grad, jit, vmap), fast and flexible, growing adoption in research. Choose PyTorch for research and prototyping (vast community, easy debugging), TensorFlow for production at scale (robust serving infrastructure), JAX for custom algorithms and performance-critical code. All three support automatic differentiation, GPU acceleration, and distributed training.",

    "Multimodal Model Libraries: HuggingFace Transformers (huggingface.co/transformers) - unified API for thousands of pre-trained models including CLIP, BLIP, LLaVA, GPT, BERT. Simple interface: model = AutoModel.from_pretrained(). Extensive documentation and tutorials. HuggingFace PEFT (huggingface.co/peft) - LoRA, adapters, prefix tuning. HuggingFace Datasets (huggingface.co/datasets) - standardized dataset loading, streaming, preprocessing. OpenCLIP (github.com/mlfoundations/open_clip) - open implementation of CLIP with training code. LAVIS (github.com/salesforce/LAVIS) - Salesforce library for vision-language tasks, includes BLIP-2. Detectron2 (github.com/facebookresearch/detectron2) - object detection and segmentation. Timm (github.com/huggingface/pytorch-image-models) - image model architectures and pre-trained weights.",

    "MLOps and Experiment Tracking: Weights & Biases (wandb.ai) - experiment tracking, hyperparameter tuning, model versioning. Beautiful visualizations, collaboration features. MLflow (mlflow.org) - open-source platform for ML lifecycle: experiment tracking, model registry, deployment. TensorBoard - visualization for TensorFlow and PyTorch training. Neptune.ai - similar to W&B, enterprise features. DVC (dvc.org) - Git for data and models, version control for datasets and experiments. Hydra (hydra.cc) - configuration management for experiments. Ray (ray.io) - distributed computing, hyperparameter tuning, serving. These tools essential for managing complex experiments and transitioning research to production.",

    "Deployment and Serving: TorchServe - PyTorch model serving with REST API, batching, metrics. Triton Inference Server (NVIDIA) - high-performance serving for PyTorch, TensorFlow, ONNX. Supports dynamic batching, model ensembles, multi-GPU. Ray Serve - Python-first serving with advanced composition. BentoML (bentoml.com) - package models as services, deploy to cloud. ONNX (onnx.ai) - open format for model interchange between frameworks. Docker for containerization. Kubernetes for orchestration at scale. Cloud platforms: AWS SageMaker, Google Vertex AI, Azure ML for managed ML infrastructure. Monitoring: Prometheus + Grafana for metrics and dashboards. These tools enable transitioning from Jupyter notebooks to production-grade systems serving millions of requests."
], img)

print("[Slide 70/78] Next Steps for Students")
img = create_simple_diagram("Next Steps", ["Lab", "Homework", "Project", "Explore"])
add_comprehensive_slide("Next Steps and Continued Learning", [
    "Immediate Actions This Week: (1) Complete the lab session - hands-on experience with CLIP, BLIP-2, and LoRA is essential for understanding concepts deeply. Experiment beyond the provided examples. (2) Start the homework assignment early - fine-tuning takes time and you may encounter technical issues. Use office hours and discussion board for questions. (3) Review lecture materials - re-read slides focusing on sections you found challenging. Mathematical formulations become clearer with multiple passes. (4) Explore the provided Jupyter notebooks - run cells, modify code, observe how changes affect results. (5) Set up your Palmetto account and test GPU access if you haven't already.",

    "Building Deeper Understanding: (1) Read the foundational papers - CLIP, BLIP-2, Flamingo, LoRA. Don't just skim - take notes, derive key equations, try to implement core ideas. (2) Reproduce paper results - pick a recent paper with released code and try to replicate their main results. This reveals implementation details omitted in papers. (3) Participate in study groups - explaining concepts to peers solidifies your own understanding and exposes knowledge gaps. (4) Follow recent research - set up arXiv alerts for your topics of interest. Multimodal AI evolves rapidly and yesterday's state-of-the-art quickly becomes baseline. (5) Attend virtual conferences if possible - NeurIPS, CVPR, ICLR often have virtual attendance options. Talks and workshops provide insights beyond papers.",

    "Project Ideas for Portfolio: (1) Domain-specific adaptation - fine-tune multimodal models for specialized domains like medical imaging, satellite imagery, or scientific diagrams. (2) Multimodal datasets - collect and curate new image-text pairs for underserved languages or domains. (3) Evaluation and analysis - comprehensive benchmarking of models on fairness, robustness, or interpretability. (4) Applications - build practical tools like accessibility applications, content moderation systems, or creative tools. (5) Efficiency - work on compression, quantization, or distillation to make models more accessible. (6) Novel architectures - experiment with new fusion strategies or attention mechanisms. Document your work through blog posts or GitHub repositories - builds your portfolio and contributes to the community.",

    "Career Development: (1) Specialize or generalize - decide whether to focus deeply on multimodal AI or maintain broader ML knowledge. Both paths valid depending on goals. (2) Build online presence - GitHub for code, personal website for portfolio, Twitter/LinkedIn for networking. (3) Contribute to open source - improves skills, visibility, and network. (4) Apply for internships and research positions - hands-on experience crucial. (5) Publish or blog - doesn't have to be top conferences. Workshop papers, arXiv preprints, and blog posts all demonstrate capability. (6) Network - attend meetups (virtual or in-person), join Discord/Slack communities, reach out to researchers whose work interests you. (7) Stay curious - the field evolves rapidly. Continuous learning essential. What you learn today provides foundation for tomorrow's advances."
], img)

print("[Slide 71/78] Course Evaluation")
img = create_simple_diagram("Feedback", ["Content", "Delivery", "Materials", "Improvement"])
add_comprehensive_slide("Course Evaluation and Feedback", [
    "Your feedback is valuable for improving this course. Please complete the official course evaluation form and consider providing feedback on: (1) Content coverage - were topics explained with appropriate depth? What would you like to see more or less of? (2) Pace - was the lecture too fast, too slow, or appropriate? (3) Materials - were slides, notebooks, and homework helpful? How can they be improved? (4) Lab and homework - were hands-on components effective for learning? Appropriate difficulty? (5) Instructor effectiveness - clear explanations, responsiveness to questions, accessibility. Specific feedback more useful than general ratings - instead of 'good lecture', explain what specifically worked or what could improve.",

    "Beyond formal evaluation, feel free to reach out with: (1) Questions on lecture material - office hours, email, or discussion board. (2) Career advice - happy to discuss research directions, internship opportunities, graduate school. (3) Research collaboration - if topics in lecture sparked project ideas. (4) Guest lectures or additional resources you'd recommend. Clemson's Applied Data Science program continually evolves and your input shapes future offerings. Previous student feedback has led to increased hands-on components, clearer mathematical explanations, and more diverse application examples.",

    "Success in this course isn't just measured by grades but by: (1) Conceptual understanding - can you explain multimodal learning principles to others? (2) Practical skills - can you implement, train, and evaluate models independently? (3) Critical thinking - can you analyze model failures, propose improvements, identify appropriate applications? (4) Research literacy - can you read papers, understand contributions, place work in broader context? (5) Ethical awareness - do you consider fairness, privacy, and societal impact when developing AI? These skills serve you throughout your career, whether in academia or industry, research or engineering, startups or established companies.",

    "Thank you for your engagement, questions, and curiosity throughout this lecture. Multimodal AI is an exciting field with profound potential for positive impact - from accessibility tools helping disabled users to scientific breakthroughs in healthcare and climate science. It also poses challenges requiring thoughtful consideration of bias, privacy, and responsible deployment. You are equipped with both technical skills and ethical frameworks to contribute to this field's future. I'm excited to see what you build and discover. Stay curious, keep learning, and don't hesitate to reach out. Good luck with the lab, homework, and your future endeavors in AI!"
], img)

print("[Slide 72/78] Additional Resources")
img = create_simple_diagram("More Resources", ["Datasets", "Models", "Tutorials", "Communities"])
add_comprehensive_slide("Additional Resources and References", [
    "Major Datasets and Benchmarks: Image-Text: MS-COCO (123K images), Flickr30K (31K), Visual Genome (108K images with detailed annotations), LAION-400M and LAION-5B (large-scale web data), Conceptual Captions (3.3M), SBU Captions (1M). Video-Text: MSR-VTT, YouCook2, ActivityNet Captions. VQA: VQAv2 (1.1M questions), GQA (compositional), OK-VQA (knowledge-intensive), VizWiz (accessibility-focused). Evaluation: ImageNet for zero-shot classification, Winoground for compositionality, MMMU for understanding, various domain-specific benchmarks. Most available via HuggingFace Datasets or official websites.",

    "Pre-trained Models: CLIP variants: OpenAI's official models (ViT-B/32, ViT-L/14), OpenCLIP trained on LAION, domain-specific versions (BiomedCLIP, GeoRSCLIP). BLIP/BLIP-2: Salesforce models with different LLM backends (OPT, FlanT5). LLaVA: 7B and 13B variants with different vision encoders. Flamingo: Limited availability (DeepMind). GPT-4V and Gemini: Proprietary via APIs. Domain-specific: medical (BiomedCLIP, Med-PaLM), scientific (Galactica), general-purpose (LLaMA, Vicuna, MPT). Model hubs: HuggingFace (main source), OpenCLIP repository, individual research project pages. Check licenses before commercial use.",

    "Interactive Demos and Playgrounds: HuggingFace Spaces - web interfaces for trying models without setup. Search for CLIP, BLIP, or specific models. Replicate.com - run models via API or web interface. Google Colab notebooks - many researchers share Colab demos of their papers. Gradio and Streamlit - frameworks for building interactive demos (good for showcasing your projects). GitHub repositories often include demo scripts. Playing with models interactively builds intuition for capabilities and limitations faster than reading papers alone.",

    "Communities and Discussion: Research communities: Reddit r/MachineLearning (general ML discussions), r/computervision, r/LanguageTechnology. Twitter/X - follow researchers, labs (@OpenAI, @MetaAI, @GoogleAI, @AnthropicAI), conference hashtags (#NeurIPS, #CVPR). Discord servers: HuggingFace, LAION, Eleuther AI, Stability AI. Forums: HuggingFace Forums for technical questions. Local: Clemson ML reading groups, AI student organizations. Conferences: NeurIPS (general ML), ICML, ICLR, CVPR (vision), ICCV, ECCV, ACL (NLP), EMNLP. Many have virtual options. Workshops at conferences often more accessible than main tracks. Networking at these venues valuable for career development and staying current with research."
], img)

print("[Slide 73/78] Looking Forward")
img = create_simple_diagram("Future", ["2024-2025", "Opportunities", "Challenges", "Impact"])
add_comprehensive_slide("Looking Forward: The Future of Multimodal AI", [
    "Near-Term Developments (2024-2025): Expect continued scaling of models and datasets, with models reaching hundreds of billions or trillions of parameters trained on tens of billions of multimodal examples. Improved architectures will reduce computational costs while increasing capabilities. Expect breakthroughs in: (1) Long-context understanding - processing hours of video or hundreds of pages of documents with images. (2) Improved reasoning - chain-of-thought prompting becoming standard, neuro-symbolic approaches gaining traction. (3) Multimodal generation - not just text-to-image but controllable video generation, 3D asset creation, interactive content. (4) Efficiency - smaller models approaching large model performance through distillation, better architectures, and training algorithms. (5) Personalization - models adapting to individual users while preserving privacy.",

    "Medium-Term Vision (2025-2030): Multimodal AI will become ubiquitous in applications. Embodied AI in robots will transform manufacturing, logistics, healthcare, and domestic assistance - robots understanding complex instructions and adapting to novel situations. Scientific research will accelerate through AI assistants that read papers, analyze data, generate hypotheses, and even design experiments. Creative tools will enable anyone to create professional-quality visual content, animations, and interactive experiences from natural language descriptions. Education will be transformed by personalized AI tutors that adapt to individual learning styles, answer questions about any content, and generate custom educational materials. Healthcare will see AI assistants helping with diagnosis, treatment planning, and patient monitoring, though human oversight will remain essential.",

    "Technical Frontiers: Key research directions will shape the field's evolution: (1) Human-AI collaboration - moving beyond AI as tool to AI as collaborator, requiring better interfaces, explainability, and mixed-initiative interaction. (2) Continual learning - models that learn and adapt over time without catastrophic forgetting, keeping pace with world changes. (3) Compositional and systematic generalization - models that truly understand compositional structure and generalize to novel combinations systematically rather than through memorization. (4) Grounding and factual accuracy - eliminating hallucination through better grounding mechanisms, retrieval augmentation, and training objectives. (5) Efficiency at scale - reducing computational and environmental costs of training and deployment while maintaining or improving capabilities. (6) Safety and alignment - ensuring AI systems behave according to human values across diverse contexts and cultures.",

    "Your Role in Shaping the Future: The multimodal AI systems of tomorrow will be built by you - this generation of students, researchers, and engineers. You'll face exciting opportunities and serious responsibilities: (1) Technical innovation - developing algorithms, architectures, and systems pushing capabilities forward. (2) Responsible development - ensuring AI benefits society broadly, mitigating harms, advancing fairness and transparency. (3) Interdisciplinary collaboration - working with domain experts, ethicists, policymakers, and affected communities. (4) Knowledge dissemination - sharing discoveries through papers, code, and educational content. (5) Mentorship - teaching the next generation. The foundations you've learned in this lecture equip you to contribute meaningfully to this field's future. We've covered theory, architectures, training, applications, and ethics - providing both technical skills and critical thinking frameworks. The field needs diverse perspectives and backgrounds to realize multimodal AI's potential while addressing its challenges. Your unique viewpoints, experiences, and ideas will shape where the field goes next."
], img)

print("[Slide 74/78] Acknowledgments")
img = create_simple_diagram("Thanks", ["Researchers", "Open Source", "Community", "Students"])
add_comprehensive_slide("Acknowledgments", [
    "This lecture builds on the work of countless researchers who have advanced multimodal AI. Special recognition to: OpenAI team for CLIP and GPT-4V, Salesforce Research for BLIP and BLIP-2, Google DeepMind for Flamingo and Gemini, University of Wisconsin-Madison team for LLaVA, Microsoft Research for LoRA, and many others whose papers we've studied. The open-source community has been instrumental - HuggingFace for democratizing access to models, LAION for creating large-scale datasets, PyTorch and TensorFlow teams for deep learning infrastructure, and countless contributors to open-source ML libraries.",

    "Clemson University's investment in computational resources through the Palmetto cluster enables hands-on learning impossible without access to GPUs. The Applied Data Science program faculty and staff support innovative curriculum development. Teaching assistants help students navigate technical challenges. Most importantly, thank you - the students - for your engagement, questions, and enthusiasm. Your curiosity and perspectives make teaching rewarding and often lead to insights I hadn't considered.",

    "The broader ML community deserves recognition for its unusual openness compared to many fields. The practice of releasing code, models, and datasets; posting preprints on arXiv; and sharing knowledge through blogs and social media accelerates progress and enables learning. While challenges around proprietary vs. open development exist, the field's commitment to openness has enabled rapid advancement and broad participation. This lecture and your learning are possible because researchers shared their work openly.",

    "Finally, acknowledgment that AI development involves complex trade-offs and that reasonable people disagree on best practices for training data, model release, safety measures, and more. The views expressed in this lecture represent my synthesis of current research and practice, but multimodal AI is an evolving field where many questions remain open. Engaging thoughtfully with these questions, considering multiple perspectives, and contributing to constructive discussions will advance the field more than any single technical innovation."
], img)

print("[Slide 75/78] References - Key Papers (1/2)")
add_comprehensive_slide("References: Key Papers (Part 1)", [
    "Foundational Architectures and Methods:\n\nVaswani, A., et al. (2017). 'Attention Is All You Need.' NeurIPS 2017. Introduces Transformer architecture.\n\nDosovitskiy, A., et al. (2020). 'An Image is Worth 16x16 Words: Transformers for Image Recognition at Scale.' ICLR 2021. Vision Transformers (ViT).\n\nRadford, A., et al. (2021). 'Learning Transferable Visual Models From Natural Language Supervision.' ICML 2021. CLIP model.\n\nJia, C., et al. (2021). 'Scaling Up Visual and Vision-Language Representation Learning With Noisy Text Supervision.' ICML 2021. ALIGN model.",

    "Advanced Multimodal Architectures:\n\nLi, J., et al. (2022). 'BLIP: Bootstrapping Language-Image Pre-training for Unified Vision-Language Understanding and Generation.' ICML 2022.\n\nLi, J., et al. (2023). 'BLIP-2: Bootstrapping Language-Image Pre-training with Frozen Image Encoders and Large Language Models.' ICML 2023.\n\nAlayrac, J., et al. (2022). 'Flamingo: a Visual Language Model for Few-Shot Learning.' NeurIPS 2022.\n\nLiu, H., et al. (2023). 'Visual Instruction Tuning.' NeurIPS 2023. LLaVA model.",

    "Training and Fine-tuning Methods:\n\nHu, E., et al. (2021). 'LoRA: Low-Rank Adaptation of Large Language Models.' ICLR 2022.\n\nOord, A., Li, Y., & Vinyals, O. (2018). 'Representation Learning with Contrastive Predictive Coding.' arXiv preprint.\n\nChen, T., et al. (2020). 'A Simple Framework for Contrastive Learning of Visual Representations.' ICML 2020. SimCLR.\n\nHe, K., et al. (2020). 'Momentum Contrast for Unsupervised Visual Representation Learning.' CVPR 2020. MoCo."
], None)

print("[Slide 76/78] References - Key Papers (2/2)")
add_comprehensive_slide("References: Key Papers (Part 2)", [
    "Evaluation and Benchmarks:\n\nLin, T., et al. (2014). 'Microsoft COCO: Common Objects in Context.' ECCV 2014.\n\nGoyal, Y., et al. (2017). 'Making the V in VQA Matter: Elevating the Role of Image Understanding in Visual Question Answering.' CVPR 2017. VQAv2.\n\nKrishna, R., et al. (2017). 'Visual Genome: Connecting Language and Vision Using Crowdsourced Dense Image Annotations.' IJCV 2017.\n\nYuksekgonul, M., et al. (2023). 'When and why vision-language models behave like bags-of-words.' ICLR 2023. Winoground.",

    "Applications and Specialized Domains:\n\nZhang, S., et al. (2023). 'BiomedCLIP: A Multimodal Biomedical Foundation Model Pretrained from Fifteen Million Scientific Image-Text Pairs.' arXiv preprint.\n\nBrohan, A., et al. (2022). 'RT-1: Robotics Transformer for Real-World Control at Scale.' arXiv preprint.\n\nRamesh, A., et al. (2021). 'Zero-Shot Text-to-Image Generation.' ICML 2021. DALL-E.\n\nRombach, R., et al. (2022). 'High-Resolution Image Synthesis with Latent Diffusion Models.' CVPR 2022. Stable Diffusion.",

    "Ethics and Responsible AI:\n\nBuolamwini, J., & Gebru, T. (2018). 'Gender Shades: Intersectional Accuracy Disparities in Commercial Gender Classification.' FAT 2018.\n\nMitchell, M., et al. (2019). 'Model Cards for Model Reporting.' FAT 2019.\n\nBender, E., et al. (2021). 'On the Dangers of Stochastic Parrots: Can Language Models Be Too Big?' FAccT 2021.\n\nGebru, T., et al. (2021). 'Datasheets for Datasets.' CACM 2021."
], None)

print("[Slide 77/78] Contact and Office Hours")
img = create_simple_diagram("Contact", ["Email", "Office Hours", "Discussion Board", "Resources"])
add_comprehensive_slide("Contact Information and Support", [
    "Instructor Contact: For questions about lecture content, homework, or course logistics, multiple channels available: (1) Office hours - check course website for schedule, both in-person and virtual options typically offered. Best for detailed technical discussions, homework help, or career advice. (2) Email - for scheduling issues, grade inquiries, or private matters. Response typically within 24-48 hours on weekdays. (3) Course discussion board - preferred for technical questions as answers benefit all students. Teaching assistants and peers often respond quickly. Check existing threads before posting new questions.",

    "Technical Support: For Palmetto cluster issues (login problems, job scheduling, quota issues), contact Clemson CITI support (support@clemson.edu) rather than course staff - they can resolve infrastructure issues faster. For HuggingFace/library issues, consult official documentation first, then discussion board. For notebook/code bugs, include error messages, environment details, and minimal reproducible example in questions. GitHub repositories for lab and homework include issue trackers for reporting bugs or suggesting improvements.",

    "Learning Resources: Course website contains all materials: slides, notebooks, homework specifications, additional readings. Updated throughout semester. Recorded lecture available on Canvas for review. Supplementary materials (papers, tutorials, dataset links) on course repository. Study groups encouraged - peer learning often more effective than studying alone. Form study groups via course discussion board or class Discord server if available.",

    "Academic Integrity: Collaboration encouraged on lab exercises and concept discussions. However, homework assignments must represent your own work. Acceptable: discussing approaches, debugging together, explaining concepts to each other. Not acceptable: copying code, sharing solutions, using solutions from previous semesters or online sources without attribution. If you use external resources (documentation, StackOverflow, ChatGPT for debugging), cite them in comments. When in doubt about whether collaboration is appropriate, ask before submitting. Academic integrity violations taken seriously and reported to university. Goal is learning - shortcuts deprive you of that. If struggling with assignments, reach out for help through appropriate channels rather than compromising integrity."
], img)

print("[Slide 78/78] Thank You and Q&A")
add_title_slide("Thank You!", "Questions and Discussion")

output_file = '/home/user/experiments/Multimodal_LLM_Lecture_COMPLETE_78_SLIDES.pptx'
prs.save(output_file)

print(f"\n{'='*80}")
print(f"✓✓✓ ALL 78 SLIDES COMPLETED! ✓✓✓")
print(f"{'='*80}")
print(f"Final presentation saved to:")
print(f"{output_file}")
print(f"Total slides: {len(prs.slides)}")
print(f"{'='*80}\n")
