#!/usr/bin/env python3
"""
Generate a comprehensive PowerPoint presentation on Multimodal Large Language Models
for a graduate-level Applied Data Science course.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Create the complete presentation on Multimodal LLMs."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme
    CLEMSON_ORANGE = RGBColor(246, 103, 51)
    CLEMSON_PURPLE = RGBColor(82, 45, 128)
    DARK_GRAY = RGBColor(51, 51, 51)
    LIGHT_GRAY = RGBColor(242, 242, 242)

    def add_title_slide(title, subtitle=""):
        """Add a title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = CLEMSON_ORANGE

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = DARK_GRAY

        return slide

    def add_content_slide(title, content_items, level=0):
        """Add a slide with title and bullet points."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = CLEMSON_PURPLE

        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]

            if isinstance(item, tuple):
                p.text = item[0]
                p.level = item[1]
            else:
                p.text = item
                p.level = level

            p.font.size = Pt(18) if p.level == 0 else Pt(16)
            p.space_before = Pt(8)
            p.font.color.rgb = DARK_GRAY

        return slide

    def add_section_slide(section_title):
        """Add a section divider slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = CLEMSON_PURPLE

        # Section title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        return slide

    # ==================== SLIDES ====================

    # Slide 1: Title
    add_title_slide(
        "Multimodal Large Language Models",
        "Applied Data Science - Clemson University"
    )

    # Slide 2: Course Context
    add_content_slide("Course Context", [
        "Part of Applied Data Science course",
        "Master's program in Computer Science",
        "Prerequisites: Understanding of Transformer architecture",
        ("Attention mechanisms", 1),
        ("Self-attention and cross-attention", 1),
        ("Encoder-decoder architectures", 1),
        "Building on previous lessons on LLMs",
        "Today: Extending to multiple modalities"
    ])

    # Slide 3: Learning Objectives
    add_content_slide("Learning Objectives", [
        "Understand multimodal learning fundamentals",
        "Master single-modality representations",
        "Learn multimodal fusion strategies",
        "Study state-of-the-art architectures (CLIP, BLIP, LLaVA, etc.)",
        "Understand training from scratch vs. fine-tuning",
        "Gain practical implementation knowledge",
        "Apply concepts on Palmetto cluster"
    ])

    # Slide 4: What is Multimodal Learning?
    add_content_slide("What is Multimodal Learning?", [
        "Definition: Learning from multiple data modalities simultaneously",
        "Modalities: Different types of data sources",
        ("Vision (images, videos)", 1),
        ("Text (documents, captions)", 1),
        ("Audio (speech, sound)", 1),
        ("Sensor data, time series", 1),
        "Goal: Build unified representations across modalities",
        "Enable cross-modal understanding and generation"
    ])

    # Slide 5: Why Multimodal Models?
    add_content_slide("Why Multimodal Models?", [
        "Human perception is inherently multimodal",
        "Richer understanding through multiple perspectives",
        "Applications benefit from combined information:",
        ("Visual Question Answering", 1),
        ("Image/Video captioning", 1),
        ("Text-to-image generation", 1),
        ("Embodied AI and robotics", 1),
        "Better generalization and robustness",
        "Alignment between different data types"
    ])

    # Slide 6: Historical Context
    add_content_slide("Historical Context", [
        "Early days: Separate models per modality",
        "2010s: Simple concatenation approaches",
        "2017: Attention mechanisms enable better fusion",
        "2019-2020: BERT + Vision models (ViLBERT, LXMERT)",
        "2021: CLIP revolutionizes vision-language learning",
        "2022-2023: Large-scale multimodal models (Flamingo, GPT-4V)",
        "2023-2024: Open-source alternatives (LLaVA, BLIP-2)",
        "2024+: Any-to-any multimodal models"
    ])

    # Slide 7: Challenges in Multimodal Learning
    add_content_slide("Challenges in Multimodal Learning", [
        "Modality gap: Different data distributions",
        "Representation learning: Finding common space",
        "Alignment: Matching corresponding elements",
        "Fusion strategies: How to combine modalities?",
        "Data requirements: Need paired/aligned data",
        "Computational cost: Processing multiple modalities",
        "Evaluation: Measuring cross-modal performance",
        "Scalability: Training with massive datasets"
    ])

    # Section: Single Modalities
    add_section_slide("Part I: Single Modality Foundations")

    # Slide 8: Single Modality - Vision
    add_content_slide("Single Modality: Vision", [
        "Image representation: Pixels arranged in grids",
        "Convolutional Neural Networks (CNNs):",
        ("Learn spatial hierarchies", 1),
        ("Translation invariance", 1),
        "Vision Transformers (ViT):",
        ("Split image into patches", 1),
        ("Apply transformer to sequence of patches", 1),
        ("Better scalability with data", 1),
        "Feature extraction: CNN backbones or ViT encoders",
        "Output: High-dimensional feature vectors"
    ])

    # Slide 9: Vision Transformers Detail
    add_content_slide("Vision Transformers (ViT)", [
        "Architecture: Pure transformer applied to images",
        "Process:",
        ("1. Split image into fixed-size patches (e.g., 16×16)", 1),
        ("2. Flatten patches and linearly embed", 1),
        ("3. Add positional embeddings", 1),
        ("4. Feed through transformer encoder", 1),
        ("5. Extract [CLS] token or average pooling", 1),
        "Advantages: Scalable, less inductive bias",
        "Used in: CLIP, BLIP, and most modern models"
    ])

    # Slide 10: Single Modality - Text
    add_content_slide("Single Modality: Text", [
        "Tokenization: Convert text to discrete tokens",
        "Embedding: Map tokens to continuous vectors",
        "Transformers for text:",
        ("BERT: Bidirectional encoder", 1),
        ("GPT: Autoregressive decoder", 1),
        ("T5: Encoder-decoder", 1),
        "Contextual representations capture semantics",
        "Pre-training on large text corpora",
        "Output: Sequence of contextual embeddings"
    ])

    # Slide 11: Single Modality - Audio
    add_content_slide("Single Modality: Audio", [
        "Raw audio: Waveform representation",
        "Feature extraction:",
        ("Mel-spectrograms: Time-frequency representation", 1),
        ("MFCCs: Mel-frequency cepstral coefficients", 1),
        "Models:",
        ("Wav2Vec 2.0: Self-supervised speech representation", 1),
        ("Whisper: Robust speech recognition", 1),
        ("AudioMAE: Masked autoencoding for audio", 1),
        "Output: Audio embeddings"
    ])

    # Slide 12: Single Modality - Video
    add_content_slide("Single Modality: Video", [
        "Video = Sequence of images + temporal information",
        "Spatial encoding: Per-frame CNN/ViT",
        "Temporal encoding:",
        ("3D convolutions", 1),
        ("Recurrent networks (LSTM)", 1),
        ("Temporal transformers", 1),
        "Challenges: Computational cost, long sequences",
        "Models: TimeSformer, VideoMAE, ViViT",
        "Output: Spatio-temporal features"
    ])

    # Slide 13: Modality Representations
    add_content_slide("Modality Representations", [
        "Each modality has unique characteristics:",
        ("Vision: Spatial structure, continuous", 1),
        ("Text: Sequential, discrete tokens", 1),
        ("Audio: Temporal, continuous waveforms", 1),
        "Goal: Map to common semantic space",
        "Representation learning objectives:",
        ("Reconstruction", 1),
        ("Contrastive learning", 1),
        ("Masked modeling", 1),
        "Enable cross-modal reasoning"
    ])

    # Slide 14: Feature Extraction Pipeline
    add_content_slide("Feature Extraction Pipeline", [
        "Input: Raw data in original modality",
        "Preprocessing:",
        ("Images: Resize, normalize, augment", 1),
        ("Text: Tokenize, add special tokens", 1),
        ("Audio: Resample, compute spectrogram", 1),
        "Encoder: Modality-specific neural network",
        "Output: Fixed or variable-length embeddings",
        "Projection: Map to common dimension",
        "Normalization: L2 normalize for similarity"
    ])

    # Slide 15: Embedding Spaces
    add_content_slide("Embedding Spaces", [
        "Embedding space: High-dimensional vector space",
        "Desirable properties:",
        ("Semantic similarity → Vector proximity", 1),
        ("Modality invariance: Same concept, same location", 1),
        ("Structured relationships preserved", 1),
        "Joint embedding space:",
        ("Multiple modalities mapped to same space", 1),
        ("Enable cross-modal retrieval", 1),
        "Learned through contrastive objectives"
    ])

    # Slide 16: The Alignment Problem
    add_content_slide("The Alignment Problem", [
        "Core challenge: How to align different modalities?",
        "Types of alignment:",
        ("Global: Image-caption pairs", 1),
        ("Local: Objects-words, regions-phrases", 1),
        ("Temporal: Video frames-transcript segments", 1),
        "Solutions:",
        ("Contrastive learning (CLIP)", 1),
        ("Cross-attention mechanisms (Flamingo)", 1),
        ("Q-Former architecture (BLIP-2)", 1),
        "Goal: Bridge the modality gap"
    ])

    # Section: Multimodal Architectures
    add_section_slide("Part II: Multimodal Architectures")

    # Slide 17: Overview of Multimodal Architectures
    add_content_slide("Overview of Multimodal Architectures", [
        "Three main paradigms:",
        "1. Dual-encoder (CLIP, ALIGN)",
        ("Separate encoders per modality", 1),
        ("Joint embedding space", 1),
        "2. Fusion encoder (LXMERT, ViLBERT)",
        ("Cross-modal attention layers", 1),
        ("Deep interaction", 1),
        "3. Encoder-decoder (BLIP, Flamingo, LLaVA)",
        ("Visual encoder + Language model", 1),
        ("Generation capabilities", 1)
    ])

    # Slide 18: Early Fusion
    add_content_slide("Early Fusion", [
        "Combine raw features at input level",
        "Concatenate or sum modality features",
        "Process jointly through unified network",
        "Advantages:",
        ("Maximum interaction between modalities", 1),
        ("Simple architecture", 1),
        "Disadvantages:",
        ("Difficult to leverage pre-trained models", 1),
        ("Modality-specific preprocessing challenging", 1),
        ("Less flexible", 1)
    ])

    # Slide 19: Late Fusion
    add_content_slide("Late Fusion", [
        "Process each modality independently",
        "Combine at decision/output level",
        "Examples: Ensemble approaches, voting",
        "Advantages:",
        ("Easy to use pre-trained models", 1),
        ("Modular design", 1),
        ("Can train modalities separately", 1),
        "Disadvantages:",
        ("Limited cross-modal interaction", 1),
        ("May miss complex relationships", 1)
    ])

    # Slide 20: Hybrid Fusion
    add_content_slide("Hybrid Fusion", [
        "Combine early and late fusion benefits",
        "Multi-stage interaction",
        "Example: Cross-attention at multiple layers",
        "Advantages:",
        ("Flexible architecture design", 1),
        ("Can leverage pre-trained components", 1),
        ("Rich cross-modal interaction", 1),
        "Most modern architectures use hybrid approach",
        "Examples: BLIP-2, Flamingo, LLaVA"
    ])

    # Slide 21: Cross-Attention Mechanisms
    add_content_slide("Cross-Attention Mechanisms", [
        "Key technique for multimodal fusion",
        "Standard attention: Q, K, V from same sequence",
        "Cross-attention: Q from one modality, K,V from another",
        "Example: Text queries vision features",
        ("Q = text embeddings", 1),
        ("K, V = image patch features", 1),
        "Bidirectional cross-attention for deep fusion",
        "Used in: Flamingo, BLIP, Q-Former"
    ])

    # Slide 22: CLIP Architecture
    add_content_slide("CLIP: Contrastive Language-Image Pre-training", [
        "Dual-encoder architecture (OpenAI, 2021)",
        "Components:",
        ("Image encoder: ViT or ResNet", 1),
        ("Text encoder: Transformer", 1),
        ("Project to same dimension", 1),
        "Training: Contrastive learning on 400M pairs",
        "Zero-shot capabilities: Natural language supervision",
        "Applications: Image classification, retrieval, generation"
    ])

    # Slide 23: CLIP Training Approach
    add_content_slide("CLIP Training Objective", [
        "Contrastive loss on batch of N pairs",
        "Maximize similarity of correct (image, text) pairs",
        "Minimize similarity of incorrect pairs",
        "Loss: InfoNCE (normalized temperature-scaled cross entropy)",
        "Math: L = -log(exp(sim(I,T)/τ) / Σ exp(sim(I,T')/τ))",
        "Symmetric: Both image→text and text→image",
        "Temperature τ: Learned parameter for scaling",
        "Result: Aligned vision-language space"
    ])

    # Slide 24: ALIGN Architecture
    add_content_slide("ALIGN: A Large-scale ImaGe and Noisy text", [
        "Similar to CLIP but different training data (Google, 2021)",
        "Key difference: Uses noisy web data (1.8B pairs)",
        "No careful curation of alt-text",
        "Components:",
        ("EfficientNet for images", 1),
        ("BERT for text", 1),
        "Same contrastive objective as CLIP",
        "Insight: Scale and diversity > data quality",
        "Achieves comparable or better results"
    ])

    # Slide 25: BLIP Architecture
    add_content_slide("BLIP: Bootstrapping Language-Image Pre-training", [
        "Encoder-decoder architecture (Salesforce, 2022)",
        "Three components:",
        ("1. Image encoder: ViT", 1),
        ("2. Text encoder: BERT-like", 1),
        ("3. Text decoder: GPT-like", 1),
        "Multi-task learning:",
        ("Image-text contrastive learning", 1),
        ("Image-text matching", 1),
        ("Image-conditioned language modeling", 1),
        "CapFilt: Bootstrapping with synthetic captions"
    ])

    # Slide 26: BLIP-2 Improvements
    add_content_slide("BLIP-2: Querying Transformer", [
        "Efficient architecture (Salesforce, 2023)",
        "Key innovation: Q-Former (Querying Transformer)",
        ("Lightweight module between frozen encoders", 1),
        ("Learnable query tokens", 1),
        ("Cross-attention to vision features", 1),
        "Two-stage training:",
        ("Stage 1: Vision-language representation learning", 1),
        ("Stage 2: Vision-to-language generative learning", 1),
        "Uses frozen ViT and frozen LLM",
        "Much more parameter-efficient than BLIP"
    ])

    # Slide 27: Flamingo Architecture
    add_content_slide("Flamingo: Visual Language Model", [
        "Few-shot learning VLM (DeepMind, 2022)",
        "Architecture:",
        ("Vision encoder: Frozen pre-trained", 1),
        ("Perceiver resampler: Compress visual features", 1),
        ("LLM: Frozen pre-trained (70B Chinchilla)", 1),
        ("Cross-attention layers: Gated XATTN-DENSE", 1),
        "Interleaved image-text input",
        "In-context learning: Few-shot prompting",
        "Powerful few-shot visual understanding"
    ])

    # Slide 28: GPT-4V Architecture Overview
    add_content_slide("GPT-4V: Multimodal GPT-4", [
        "Proprietary multimodal model (OpenAI, 2023)",
        "Accepts both images and text as input",
        "Architecture details not fully disclosed",
        "Likely components:",
        ("Vision encoder: Enhanced CLIP-like", 1),
        ("Large language model: GPT-4", 1),
        ("Sophisticated fusion mechanism", 1),
        "Capabilities: VQA, OCR, reasoning, code generation",
        "State-of-the-art performance on benchmarks"
    ])

    # Slide 29: LLaVA Architecture
    add_content_slide("LLaVA: Large Language and Vision Assistant", [
        "Open-source VLM (Liu et al., 2023)",
        "Simple and effective architecture:",
        ("Vision encoder: CLIP ViT-L/14", 1),
        ("Projection layer: Linear or MLP", 1),
        ("LLM: Vicuna (LLaMA-based)", 1),
        "Training data: GPT-4 generated instructions",
        "Two-stage training:",
        ("Pre-training: Projection layer only", 1),
        ("Fine-tuning: Full model or LoRA", 1),
        "Excellent performance with modest compute"
    ])

    # Slide 30: LLaVA 1.5 Enhancements
    add_content_slide("LLaVA 1.5 Improvements", [
        "Enhanced version (2023)",
        "Improvements:",
        ("MLP projection layer: 2-layer vs 1-layer", 1),
        ("Higher resolution: 336×336 instead of 224×224", 1),
        ("Better instruction data: More diverse tasks", 1),
        ("Academic task data: VQA, OCR, reasoning", 1),
        "Unified training recipe",
        "Near GPT-4V performance on some tasks",
        "Efficient: Can fine-tune on single GPU"
    ])

    # Slide 31: Gemini Architecture
    add_content_slide("Gemini: Google's Multimodal Model", [
        "Native multimodal model (Google, 2023)",
        "Three sizes: Nano, Pro, Ultra",
        "Key feature: Natively multimodal from scratch",
        ("Not adapting separate pre-trained models", 1),
        ("Joint training on all modalities", 1),
        "Capabilities:",
        ("Image, video, audio, text understanding", 1),
        ("Long context: Up to 1M tokens", 1),
        ("Complex reasoning across modalities", 1),
        "State-of-the-art on many benchmarks"
    ])

    # Slide 32: Any-to-Any Models
    add_content_slide("Any-to-Any Multimodal Models", [
        "Next generation: Unified models for all modalities",
        "Examples:",
        ("NExT-GPT: Any-to-any generation", 1),
        ("CoDi: Composable diffusion", 1),
        ("Unified-IO 2: Single model, any task", 1),
        "Capabilities:",
        ("Text → Image, Audio, Video", 1),
        ("Image → Text, Audio", 1),
        ("Any modality → Any modality", 1),
        "Challenges: Training complexity, data requirements"
    ])

    # Slide 33: Unified Multimodal Representations
    add_content_slide("Unified Multimodal Models", [
        "Goal: Single model handling all modalities",
        "Approaches:",
        ("Tokenize all modalities (Unified-IO)", 1),
        ("Shared encoder-decoder", 1),
        ("Modality-specific adapters", 1),
        "Benefits:",
        ("Simplified deployment", 1),
        ("Cross-modal transfer learning", 1),
        ("More efficient inference", 1),
        "Research direction: Foundation models"
    ])

    # Slide 34: Architecture Comparison
    add_content_slide("Architecture Comparison", [
        "CLIP: Fast, dual-encoder, zero-shot classification",
        "BLIP: Versatile, generation + understanding",
        "BLIP-2: Efficient, frozen components, Q-Former",
        "Flamingo: Few-shot, interleaved input, large-scale",
        "LLaVA: Open-source, instruction-tuned, accessible",
        "GPT-4V: State-of-the-art, proprietary, expensive",
        "Selection criteria: Task, compute, data, performance"
    ])

    # Slide 35: Architecture Design Principles
    add_content_slide("Design Principles", [
        "1. Leverage pre-trained models when possible",
        "2. Freeze large components to reduce compute",
        "3. Use efficient fusion mechanisms (Q-Former, projectors)",
        "4. Multi-task learning for better generalization",
        "5. High-quality instruction data crucial",
        "6. Balance model size vs. performance",
        "7. Consider inference efficiency",
        "8. Open-source enables research and applications"
    ])

    # Section: Training Strategies
    add_section_slide("Part III: Training Strategies")

    # Slide 36: Training from Scratch Overview
    add_content_slide("Training from Scratch", [
        "When to train from scratch?",
        ("Novel architectures", 1),
        ("Unique data distribution", 1),
        ("Proprietary requirements", 1),
        "Challenges:",
        ("Massive data requirements (millions of pairs)", 1),
        ("Extensive compute (thousands of GPU hours)", 1),
        ("Careful hyperparameter tuning", 1),
        "Alternative: Fine-tune pre-trained models",
        "Most practitioners: Fine-tuning approach"
    ])

    # Slide 37: Data Requirements
    add_content_slide("Data Requirements", [
        "Scale: Modern models trained on millions/billions of pairs",
        "Types of data:",
        ("Image-caption pairs", 1),
        ("Interleaved image-text documents", 1),
        ("Video-subtitle pairs", 1),
        ("Question-answer pairs", 1),
        "Data quality vs. quantity:",
        ("CLIP: Curated 400M pairs", 1),
        ("ALIGN: Noisy 1.8B pairs", 1),
        "Diversity crucial for generalization"
    ])

    # Slide 38: Dataset Examples
    add_content_slide("Key Datasets", [
        "Image-Text:",
        ("COCO: 330K images with captions", 1),
        ("Visual Genome: 108K images with dense annotations", 1),
        ("Conceptual Captions: 3.3M image-caption pairs", 1),
        ("LAION-400M/5B: Large-scale web data", 1),
        "VQA:",
        ("VQAv2: 1.1M questions", 1),
        ("GQA: 22M questions with scene graphs", 1),
        "Video: WebVid, HowTo100M, YT-Temporal"
    ])

    # Slide 39: Pre-training Objectives
    add_content_slide("Pre-training Objectives", [
        "1. Contrastive Learning:",
        ("Align positive pairs, separate negative pairs", 1),
        ("InfoNCE loss", 1),
        "2. Masked Language/Image Modeling:",
        ("Mask tokens/patches, predict them", 1),
        ("Learn contextualized representations", 1),
        "3. Image-Text Matching:",
        ("Binary classification: match or not", 1),
        ("Cross-modal alignment", 1),
        "4. Generative objectives:",
        ("Caption generation, image generation", 1)
    ])

    # Slide 40: Contrastive Learning Deep Dive
    add_content_slide("Contrastive Learning", [
        "Core idea: Similar items close, dissimilar far apart",
        "For multimodal: Match corresponding modalities",
        "Positive pairs: (image, caption) from same example",
        "Negative pairs: All other combinations in batch",
        "Loss function: InfoNCE",
        ("Numerator: Similarity of positive pair", 1),
        ("Denominator: Sum of all pair similarities", 1),
        "Temperature scaling: Control distribution sharpness",
        "Effective for learning aligned embeddings"
    ])

    # Slide 41: Masked Modeling
    add_content_slide("Masked Language/Image Modeling", [
        "Masked Language Modeling (MLM):",
        ("Mask random tokens in text", 1),
        ("Predict masked tokens from context", 1),
        ("Used in BERT, VilBERT", 1),
        "Masked Image Modeling (MIM):",
        ("Mask random patches in image", 1),
        ("Reconstruct masked patches", 1),
        ("Used in MAE, SimMIM", 1),
        "Multimodal: Mask in both modalities",
        "Learn robust cross-modal representations"
    ])

    # Slide 42: Image-Text Matching
    add_content_slide("Image-Text Matching (ITM)", [
        "Binary classification task",
        "Input: Image and text pair",
        "Output: Match (1) or not match (0)",
        "Requires cross-modal interaction",
        ("Usually through cross-attention", 1),
        "Negative sampling strategies:",
        ("Random negatives from batch", 1),
        ("Hard negatives: Similar but wrong", 1),
        "Used alongside contrastive learning",
        "Improves fine-grained alignment"
    ])

    # Slide 43: Training Stability
    add_content_slide("Training Stability", [
        "Challenges:",
        ("Large batch sizes needed for contrastive learning", 1),
        ("Gradient explosion/vanishing", 1),
        ("Modality imbalance", 1),
        "Solutions:",
        ("Gradient clipping", 1),
        ("Layer normalization", 1),
        ("Warm-up learning rate schedule", 1),
        ("Mixed precision training (FP16/BF16)", 1),
        ("Careful initialization", 1),
        "Monitor: Loss curves, gradient norms, representations"
    ])

    # Slide 44: Computational Requirements
    add_content_slide("Computational Requirements", [
        "Training large multimodal models is expensive",
        "Examples:",
        ("CLIP: 592 V100 GPUs, 12 days", 1),
        ("Flamingo: TPUv4 pods, weeks", 1),
        ("LLaVA: 8 A100 GPUs, ~1 day", 1),
        "Factors affecting cost:",
        ("Model size, batch size, data size, resolution", 1),
        "Strategies to reduce cost:",
        ("Freeze components, efficient architectures", 1),
        ("Lower resolution during pre-training", 1),
        ("Use smaller models for research", 1)
    ])

    # Slide 45: Distributed Training
    add_content_slide("Distributed Training", [
        "Essential for large-scale training",
        "Strategies:",
        ("Data parallelism: Split data across GPUs", 1),
        ("Model parallelism: Split model across GPUs", 1),
        ("Pipeline parallelism: Split layers across GPUs", 1),
        ("Tensor parallelism: Split tensors across GPUs", 1),
        "Frameworks:",
        ("PyTorch DDP, FSDP", 1),
        ("DeepSpeed, Megatron-LM", 1),
        "Communication overhead: Critical bottleneck"
    ])

    # Slide 46: Training Pipeline
    add_content_slide("Training Pipeline", [
        "1. Data preparation:",
        ("Collect and clean data", 1),
        ("Create train/val/test splits", 1),
        "2. Preprocessing:",
        ("Tokenization, image resizing", 1),
        ("Augmentation strategies", 1),
        "3. Model initialization:",
        ("Load pre-trained weights if available", 1),
        "4. Training loop:",
        ("Forward pass, loss computation", 1),
        ("Backward pass, optimizer step", 1),
        "5. Evaluation and checkpointing"
    ])

    # Slide 47: Evaluation Metrics
    add_content_slide("Evaluation Metrics", [
        "Image-Text Retrieval:",
        ("Recall@K, Mean Rank", 1),
        "VQA:",
        ("Accuracy, VQA score", 1),
        "Image Captioning:",
        ("BLEU, METEOR, CIDEr, SPICE", 1),
        "Zero-shot Classification:",
        ("Top-1/Top-5 accuracy", 1),
        "General:",
        ("Cross-modal similarity scores", 1),
        ("Human evaluation for generation tasks", 1)
    ])

    # Slide 48: Benchmarks
    add_content_slide("Key Benchmarks", [
        "Classification: ImageNet, CIFAR",
        "VQA: VQAv2, GQA, OK-VQA, VizWiz",
        "Captioning: COCO Captions, NoCaps",
        "Retrieval: COCO, Flickr30K",
        "Reasoning: NLVR2, CLEVR",
        "OCR: TextVQA, DocVQA, Infographics",
        "Video: MSRVTT, ActivityNet",
        "Comprehensive: MMBench, SEED-Bench, MM-Vet"
    ])

    # Section: Fine-tuning
    add_section_slide("Part IV: Fine-tuning Strategies")

    # Slide 49: Fine-tuning Overview
    add_content_slide("Fine-tuning Overview", [
        "Most practical approach for practitioners",
        "Start with pre-trained multimodal model",
        "Adapt to specific task or domain",
        "Advantages:",
        ("Less data required", 1),
        ("Less compute required", 1),
        ("Faster convergence", 1),
        ("Better performance on downstream tasks", 1),
        "Types: Full fine-tuning vs. parameter-efficient"
    ])

    # Slide 50: Transfer Learning for Multimodal
    add_content_slide("Transfer Learning", [
        "Leverage knowledge from pre-training",
        "Pre-trained model already understands:",
        ("Visual concepts", 1),
        ("Language semantics", 1),
        ("Cross-modal alignment", 1),
        "Fine-tuning teaches:",
        ("Task-specific patterns", 1),
        ("Domain-specific knowledge", 1),
        ("Output format requirements", 1),
        "Critical: Choose appropriate base model"
    ])

    # Slide 51: Parameter-Efficient Fine-tuning (PEFT)
    add_content_slide("Parameter-Efficient Fine-tuning (PEFT)", [
        "Goal: Fine-tune with minimal trainable parameters",
        "Why PEFT?",
        ("Large models too expensive to fully fine-tune", 1),
        ("Reduced memory requirements", 1),
        ("Faster training", 1),
        ("Easier deployment (store only adapter weights)", 1),
        "Methods:",
        ("LoRA, Adapters, Prompt tuning, Prefix tuning", 1),
        "Often achieves comparable performance"
    ])

    # Slide 52: LoRA for Multimodal Models
    add_content_slide("LoRA: Low-Rank Adaptation", [
        "Key idea: Learn low-rank updates to weight matrices",
        "W_new = W_frozen + B × A",
        ("W_frozen: Original pre-trained weights", 1),
        ("B, A: Low-rank matrices (rank r << d)", 1),
        "Only train B and A (1-2% of parameters)",
        "Apply to:",
        ("Attention projection matrices (Q, K, V)", 1),
        ("Feed-forward layers", 1),
        "Works well for LLM component in VLMs",
        "Implementation: PEFT library by HuggingFace"
    ])

    # Slide 53: Adapter Modules
    add_content_slide("Adapter Modules", [
        "Small bottleneck layers inserted in network",
        "Structure:",
        ("Down-projection: d → r", 1),
        ("Non-linearity", 1),
        ("Up-projection: r → d", 1),
        ("Residual connection", 1),
        "Inserted after attention and FFN layers",
        "Only adapters are trained",
        "Slightly more parameters than LoRA",
        "Good for vision encoders"
    ])

    # Slide 54: Prompt Tuning
    add_content_slide("Prompt Tuning", [
        "Learn soft prompts: Continuous embeddings",
        "Prepend learnable tokens to input",
        "Visual prompt tuning:",
        ("Learn perturbations to input image", 1),
        ("Or learn visual prefix tokens", 1),
        "Text prompt tuning:",
        ("Learn continuous prompt embeddings", 1),
        "Extremely parameter-efficient",
        "Challenges: May need more careful tuning",
        "Good for: Few-shot learning scenarios"
    ])

    # Slide 55: Instruction Tuning
    add_content_slide("Instruction Tuning", [
        "Fine-tune on instruction-following data",
        "Format: (Instruction, Input, Output) triples",
        "Example:",
        ("Instruction: 'Describe the image in detail'", 1),
        ("Input: [Image]", 1),
        ("Output: 'A photo of...'", 1),
        "Benefits:",
        ("Better instruction following", 1),
        ("More natural interactions", 1),
        ("Improved zero-shot generalization", 1),
        "Used in: LLaVA, InstructBLIP, Otter"
    ])

    # Slide 56: Task-Specific Fine-tuning
    add_content_slide("Task-Specific Fine-tuning", [
        "Adapt to specific downstream tasks",
        "Examples:",
        ("VQA: Fine-tune on VQAv2", 1),
        ("Captioning: Fine-tune on COCO", 1),
        ("Medical imaging: Fine-tune on medical data", 1),
        "Task-specific head may be added",
        "Data: 1K-100K task-specific examples",
        "Often use full fine-tuning or LoRA",
        "Evaluation on task benchmark"
    ])

    # Slide 57: Domain Adaptation
    add_content_slide("Domain Adaptation", [
        "Adapt to specific domain:",
        ("Medical images and reports", 1),
        ("Scientific figures and papers", 1),
        ("E-commerce products and descriptions", 1),
        ("Remote sensing imagery", 1),
        "Challenges:",
        ("Domain shift from pre-training data", 1),
        ("Specialized terminology", 1),
        "Solutions:",
        ("Continue pre-training on domain data", 1),
        ("Task-specific fine-tuning", 1),
        ("Domain-specific vocabularies", 1)
    ])

    # Slide 58: Fine-tuning Best Practices
    add_content_slide("Fine-tuning Best Practices", [
        "1. Start with appropriate pre-trained model",
        "2. Use lower learning rate than pre-training",
        "3. Freeze vision encoder initially (if needed)",
        "4. Use LoRA for large models",
        "5. Monitor validation performance carefully",
        "6. Prevent overfitting:",
        ("Early stopping, regularization", 1),
        "7. Data augmentation for small datasets",
        "8. Gradient accumulation for larger batch sizes"
    ])

    # Section: Implementation
    add_section_slide("Part V: Implementation & Applications")

    # Slide 59: Implementation Frameworks
    add_content_slide("Implementation Frameworks", [
        "HuggingFace Transformers:",
        ("Pre-trained models: CLIP, BLIP, LLaVA", 1),
        ("Easy fine-tuning with Trainer API", 1),
        ("PEFT library for LoRA", 1),
        "PyTorch: Low-level flexibility",
        "JAX/Flax: XLA compilation, TPU support",
        "LangChain: Integration with applications",
        "OpenCLIP: Open-source CLIP training",
        "LLaVA codebase: Instruction tuning pipeline"
    ])

    # Slide 60: Hardware Requirements
    add_content_slide("Hardware Requirements", [
        "GPU memory critical bottleneck",
        "Model size considerations:",
        ("CLIP ViT-B: ~400MB, inference on any GPU", 1),
        ("BLIP-2: ~4GB, fine-tune on 16GB+ GPU", 1),
        ("LLaVA-7B: ~14GB, fine-tune on 24GB+ GPU", 1),
        ("LLaVA-13B: ~26GB, need 40GB+ GPU", 1),
        "Techniques to reduce memory:",
        ("8-bit quantization, gradient checkpointing", 1),
        ("LoRA reduces memory by 2-3x", 1),
        "Palmetto: A100 (40/80GB), H100 available"
    ])

    # Slide 61: Palmetto Cluster Setup
    add_content_slide("Palmetto Cluster Setup", [
        "Clemson's HPC cluster",
        "GPU nodes: V100, A100, H100",
        "Setup steps:",
        ("1. Request GPU node: qsub -I -l select=1:ngpus=1", 1),
        ("2. Load modules: CUDA, conda", 1),
        ("3. Create environment: conda create -n multimodal", 1),
        ("4. Install packages: transformers, torch, peft", 1),
        "Storage: /scratch for datasets",
        "Jupyter notebooks available via OnDemand",
        "See lab notebook for detailed setup"
    ])

    # Slide 62: Code Walkthrough - Loading Models
    add_content_slide("Loading Pre-trained Models", [
        "HuggingFace example:",
        "",
        "from transformers import CLIPModel, CLIPProcessor",
        "model = CLIPModel.from_pretrained('openai/clip-vit-base-patch32')",
        "processor = CLIPProcessor.from_pretrained('openai/clip-vit-base-patch32')",
        "",
        "For BLIP-2:",
        "from transformers import Blip2Processor, Blip2ForConditionalGeneration",
        "processor = Blip2Processor.from_pretrained('Salesforce/blip2-opt-2.7b')",
        "model = Blip2ForConditionalGeneration.from_pretrained('Salesforce/blip2-opt-2.7b')",
        "",
        "Move to GPU: model.to('cuda')"
    ])

    # Slide 63: Code Walkthrough - Inference
    add_content_slide("Inference Example", [
        "CLIP zero-shot classification:",
        "",
        "from PIL import Image",
        "image = Image.open('photo.jpg')",
        "text = ['a photo of a cat', 'a photo of a dog']",
        "",
        "inputs = processor(text=text, images=image, return_tensors='pt', padding=True)",
        "outputs = model(**inputs)",
        "logits_per_image = outputs.logits_per_image",
        "probs = logits_per_image.softmax(dim=1)",
        "",
        "See lab notebook for BLIP captioning, LLaVA VQA examples"
    ])

    # Slide 64: Code Walkthrough - Fine-tuning
    add_content_slide("Fine-tuning Example", [
        "Using HuggingFace Trainer:",
        "",
        "from transformers import Trainer, TrainingArguments",
        "from peft import get_peft_model, LoraConfig",
        "",
        "lora_config = LoraConfig(r=8, lora_alpha=16, target_modules=['q_proj', 'v_proj'])",
        "model = get_peft_model(model, lora_config)",
        "",
        "training_args = TrainingArguments(",
        "    output_dir='./results', num_train_epochs=3,",
        "    per_device_train_batch_size=4, learning_rate=2e-5)",
        "",
        "trainer = Trainer(model=model, args=training_args, train_dataset=dataset)",
        "trainer.train()"
    ])

    # Slide 65: Application - Visual Question Answering
    add_content_slide("Application: Visual Question Answering", [
        "Task: Answer questions about images",
        "Input: Image + Question",
        "Output: Answer (text)",
        "Models: BLIP-2, LLaVA, InstructBLIP",
        "Use cases:",
        ("Accessibility for visually impaired", 1),
        ("Content moderation", 1),
        ("Educational tools", 1),
        ("Image search and retrieval", 1),
        "Datasets: VQAv2, GQA, OK-VQA"
    ])

    # Slide 66: Application - Image Captioning
    add_content_slide("Application: Image Captioning", [
        "Task: Generate textual description of images",
        "Input: Image",
        "Output: Caption (text)",
        "Models: BLIP, BLIP-2, GIT",
        "Use cases:",
        ("Alt-text generation for web accessibility", 1),
        ("Photo organization and search", 1),
        ("Social media content creation", 1),
        ("Medical report generation", 1),
        "Evaluation: BLEU, CIDEr, SPICE metrics"
    ])

    # Slide 67: Application - Text-to-Image Generation
    add_content_slide("Application: Text-to-Image Generation", [
        "Task: Generate images from text descriptions",
        "Input: Text prompt",
        "Output: Image",
        "Models: DALL-E, Stable Diffusion, Imagen",
        "Architecture: Diffusion models with text conditioning",
        ("Text encoder: CLIP or T5", 1),
        ("Diffusion model: U-Net with cross-attention", 1),
        "Use cases:",
        ("Creative content generation", 1),
        ("Design and prototyping", 1),
        ("Data augmentation", 1)
    ])

    # Slide 68: Application - Video Understanding
    add_content_slide("Application: Video Understanding", [
        "Task: Understand and reason about video content",
        "Input: Video (sequence of frames + audio)",
        "Output: Caption, answer, action recognition",
        "Models: VideoMAE, TimeSformer, Video-LLaMA",
        "Challenges:",
        ("Temporal modeling", 1),
        ("Long sequences (memory)", 1),
        ("Multi-modal fusion (visual + audio + text)", 1),
        "Applications: Surveillance, sports analysis, education"
    ])

    # Slide 69: Challenges and Limitations
    add_content_slide("Challenges and Limitations", [
        "Hallucination: Generating false information",
        "Robustness: Sensitive to input variations",
        "Bias: Reflecting dataset biases",
        "Computational cost: Training and inference",
        "Data requirements: Need for large paired datasets",
        "Interpretability: Understanding model decisions",
        "Safety: Potential for misuse",
        "Evaluation: Measuring true understanding"
    ])

    # Slide 70: Future Directions
    add_content_slide("Future Directions", [
        "Scaling: Larger models, more data, more modalities",
        "Efficiency: Smaller models with better performance",
        "Grounding: Connecting to physical world",
        "Reasoning: Better logical and causal reasoning",
        "Embodied AI: Integration with robotics",
        "Multilinguality: Support for more languages",
        "Scientific applications: Discovery and analysis",
        "Democratization: Making models more accessible"
    ])

    # Slide 71: Key Takeaways
    add_content_slide("Key Takeaways", [
        "Multimodal learning enables richer AI understanding",
        "Modern architectures leverage pre-trained components",
        "Contrastive learning is key for alignment",
        "Fine-tuning is practical approach for most applications",
        "Parameter-efficient methods (LoRA) make fine-tuning accessible",
        "Many open-source models available (CLIP, BLIP, LLaVA)",
        "Applications span from VQA to text-to-image generation",
        "Exciting field with rapid progress"
    ])

    # Slide 72: Resources
    add_content_slide("Resources", [
        "Papers:",
        ("CLIP, BLIP-2, LLaVA, Flamingo", 1),
        "Code:",
        ("HuggingFace Transformers, PEFT library", 1),
        ("OpenCLIP, LLaVA GitHub", 1),
        "Datasets:",
        ("COCO, Visual Genome, LAION", 1),
        "Tutorials:",
        ("HuggingFace documentation", 1),
        ("Papers with Code", 1),
        "Next: Lab session with hands-on implementation!"
    ])

    # Slide 73: Questions & Discussion
    add_title_slide("Questions & Discussion", "")

    return prs

def main():
    """Main function to generate the presentation."""
    print("Generating Multimodal LLM presentation...")
    prs = create_presentation()

    output_path = "/home/user/experiments/Multimodal_LLM_Lecture.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
