#!/usr/bin/env python3
"""Add slide 57"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-56.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

print("[Slide 57/78] Adding: Production Best Practices")

img = create_simple_diagram("Best Practices", ["Testing", "Monitoring", "Versioning", "Documentation"])
add_comprehensive_slide("Production Best Practices", [
    "Testing and Validation: Rigorous testing prevents production failures. Testing hierarchy: (1) Unit tests - test individual components (data preprocessing, model forward pass, postprocessing) with fixtures. Use pytest for Python testing. (2) Integration tests - test end-to-end pipeline: input → preprocessing → inference → postprocessing → output. Validate on representative examples covering common cases and edge cases. (3) Regression tests - maintain test suite of examples with expected outputs. Run on every code change to catch regressions. (4) Performance tests - measure latency and throughput under load using tools like Locust or Apache JMeter. Establish performance baselines and SLOs (Service Level Objectives). (5) Shadow deployment - run new model alongside production model, compare outputs, monitor disagreements before switching traffic. (6) A/B testing - gradually roll out new model to subset of users, compare metrics (accuracy, user engagement, business KPIs) before full deployment.",

    "Monitoring and Observability: Comprehensive monitoring enables quick issue detection and diagnosis. Key metrics to track: (1) Latency - p50, p95, p99, p99.9 response times. Set alerts for degradation (e.g., p95 > 200ms). (2) Throughput - requests per second, ensuring capacity matches demand. (3) Error rates - 4xx (client errors), 5xx (server errors), model-specific errors (invalid inputs). (4) Resource utilization - CPU, GPU, memory usage. High utilization may indicate need to scale. (5) Model metrics - prediction confidence distribution, embedding norms, attention pattern statistics. Detect distribution shift. (6) Business metrics - task success rates, user engagement. Infrastructure: Use Prometheus for metrics collection, Grafana for dashboards, alerting systems (PagerDuty, Opsgenie) for on-call. Structured logging with context (request ID, user ID, model version) enables debugging. Distributed tracing (Jaeger, Zipkin) for multi-service calls.",

    "Model Versioning and Reproducibility: Systematic versioning prevents confusion and enables rollback. Version control: (1) Code versioning - Git for all code, including training scripts, preprocessing, evaluation. Tag releases (v1.0, v1.1). (2) Data versioning - DVC (Data Version Control) or similar tools track dataset versions. Critical for reproducibility - 'trained on dataset version X'. (3) Model versioning - semantic versioning (major.minor.patch). Store model weights, configuration, hyperparameters, and training metadata together. Use MLflow, Weights & Biases, or Neptune for experiment tracking. (4) Environment versioning - Docker images with pinned dependency versions ensure consistent execution across development and production. (5) Experiment tracking - log every training run with hyperparameters, metrics, artifacts. Enable comparing experiments and reproducing best models. Model registry - centralized repository (MLflow Model Registry, SageMaker Model Registry) tracks production-ready models with metadata, lineage, and approval workflows.",

    "Documentation and Knowledge Sharing: Good documentation enables maintenance and knowledge transfer. Documentation layers: (1) Code documentation - docstrings for functions/classes explaining inputs, outputs, behavior. Follow Google or NumPy docstring conventions. (2) Model cards - structured documentation template describing model architecture, training data, intended use, limitations, evaluation results, ethical considerations. Essential for responsible AI deployment. (3) API documentation - clear specification of endpoints, input/output formats, error codes, rate limits. Use OpenAPI/Swagger for REST APIs. (4) Deployment documentation - runbooks for deployment procedures, rollback procedures, troubleshooting common issues, on-call playbooks. (5) Architecture documentation - system diagrams showing components, data flow, dependencies. (6) Decision logs - ADRs (Architecture Decision Records) documenting key technical decisions and rationale. Maintenance: Documentation must be living - update with code changes. Automated doc generation (Sphinx, pdoc) from code comments. Regular reviews to keep current. Knowledge sharing: internal tech talks, post-mortems for incidents, onboarding guides for new team members."
], img)

output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-57.pptx'
prs.save(output_file)
print(f"✓ Slide 57 completed: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
