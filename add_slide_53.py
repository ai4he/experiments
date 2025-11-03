#!/usr/bin/env python3
"""Add slide 53"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-52.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

print("[Slide 53/78] Adding: BLIP-2 Image Captioning and VQA Exercise")

img = create_simple_diagram("BLIP-2 Exercise", ["Load\nBLIP-2", "Caption\nImages", "Ask\nQuestions", "Analyze"])
add_comprehensive_slide("Lab Exercise 2: BLIP-2 Image Captioning and VQA", [
    "Exercise Overview: In this exercise, you'll use BLIP-2 for two tasks: image captioning (generating textual descriptions) and visual question answering (answering natural language questions about images). BLIP-2's Q-Former architecture efficiently connects a frozen vision encoder to a large language model, enabling both understanding and generation tasks. You'll experiment with different LLM backends (OPT-2.7B, FlanT5-XL) and compare their outputs. The exercise demonstrates how multimodal models can generate fluent, contextually appropriate language grounded in visual content. Estimated time: 20 minutes. Models to use: Salesforce/blip2-opt-2.7b and Salesforce/blip2-flan-t5-xl.",

    "Image Captioning Implementation: Load BLIP-2 for captioning: from transformers import Blip2Processor, Blip2ForConditionalGeneration; processor = Blip2Processor.from_pretrained('Salesforce/blip2-opt-2.7b'); model = Blip2ForConditionalGeneration.from_pretrained('Salesforce/blip2-opt-2.7b', torch_dtype=torch.float16). Move to GPU. For each test image: (1) Preprocess: inputs = processor(images=image, return_tensors='pt').to('cuda', torch.float16). (2) Generate caption: generated_ids = model.generate(**inputs, max_length=50). (3) Decode: caption = processor.batch_decode(generated_ids, skip_special_tokens=True)[0].strip(). (4) Display image with generated caption. Try multiple images - portraits, landscapes, indoor/outdoor scenes, complex multi-object scenes. Observe caption quality: Are they factually accurate? Do they capture salient details? Generic or specific?",

    "Visual Question Answering Implementation: For VQA, extend the captioning code to accept text prompts: prompt = 'Question: What color is the sky? Answer:'; inputs = processor(images=image, text=prompt, return_tensors='pt').to('cuda', torch.float16); out = model.generate(**inputs, max_length=20); answer = processor.decode(out[0], skip_special_tokens=True). Experiment with different question types: (1) Object recognition: 'What objects are in this image?', 'Is there a dog in this image?'. (2) Attribute questions: 'What color is the car?', 'How many people are there?'. (3) Spatial questions: 'Where is the cat?', 'What is behind the tree?'. (4) Activity questions: 'What is the person doing?'. (5) Reasoning questions: 'Why is the person smiling?', 'What time of day is it?'. Document which types BLIP-2 handles well vs. struggles with.",

    "Comparing LLM Backends and Analysis: Compare BLIP-2 with OPT-2.7B vs. FlanT5-XL backends on the same images and questions. Observations to make: (1) Caption quality differences - which produces more detailed, accurate, or natural descriptions? (2) Answer format - does one backend produce more concise answers? (3) Hallucination - do models describe objects not present? (4) Factual accuracy vs. fluency trade-off. (5) Response time and memory usage differences (FlanT5-XL is larger, ~4GB vs 3GB). Discussion questions: How does BLIP-2 compare to CLIP for zero-shot tasks? What are the trade-offs between the two models? When would you choose BLIP-2 over CLIP or vice versa? Can you identify failure cases - images or questions where BLIP-2 fails? How might these models be improved? This analysis builds intuition for selecting appropriate models for different applications."
], img)

output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-53.pptx'
prs.save(output_file)
print(f"✓ Slide 53 completed: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
