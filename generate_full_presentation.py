#!/usr/bin/env python3
"""
Generate complete Multimodal LLM presentation with images on every slide.
Status updates printed for each slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
import io
import sys

# Counter for status updates
slide_count = 0

def log_status(message):
    """Print status update for each slide."""
    global slide_count
    slide_count += 1
    print(f"[Slide {slide_count:2d}] {message}")
    sys.stdout.flush()

# All diagram creation functions
def create_diagram_multimodal_learning():
    """Multimodal learning concept diagram."""
    fig, ax = plt.subplots(figsize=(10, 6))
    modalities = ['Vision\n(Images)', 'Text\n(Language)', 'Audio\n(Speech)', 'Video\n(Motion)']
    colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A']
    center_x, center_y = 0.5, 0.5
    radius = 0.3
    angles = [0, 90, 180, 270]

    for i, (modality, color, angle) in enumerate(zip(modalities, colors, angles)):
        x = center_x + radius * np.cos(np.radians(angle))
        y = center_y + radius * np.sin(np.radians(angle))
        box = patches.FancyBboxPatch((x-0.08, y-0.05), 0.16, 0.1, boxstyle="round,pad=0.01",
                                     edgecolor=color, facecolor=color, alpha=0.3, linewidth=2)
        ax.add_patch(box)
        ax.text(x, y, modality, ha='center', va='center', fontsize=10, fontweight='bold')
        ax.annotate('', xy=(center_x, center_y), xytext=(x, y),
                   arrowprops=dict(arrowstyle='->', lw=2, color=color))

    center_circle = patches.Circle((center_x, center_y), 0.08, edgecolor='#2E86AB',
                                  facecolor='#A23B72', linewidth=3, alpha=0.6)
    ax.add_patch(center_circle)
    ax.text(center_x, center_y, 'Unified\nSpace', ha='center', va='center',
           fontsize=9, fontweight='bold', color='white')
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    ax.set_title('Multimodal Learning', fontsize=13, fontweight='bold')

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_simple_diagram(title, boxes, arrows=None):
    """Create a simple flow diagram."""
    fig, ax = plt.subplots(figsize=(9, 5))

    colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DFE6E9', '#6C5CE7']

    x_step = 0.8 / len(boxes)
    x_start = 0.1

    for i, box_text in enumerate(boxes):
        x = x_start + i * x_step
        color = colors[i % len(colors)]
        rect = patches.FancyBboxPatch((x, 0.4), x_step * 0.85, 0.2,
                                     boxstyle="round,pad=0.01",
                                     edgecolor='black', facecolor=color,
                                     alpha=0.5, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + x_step * 0.425, 0.5, box_text, ha='center', va='center',
               fontsize=9, fontweight='bold', wrap=True)

        if i < len(boxes) - 1:
            ax.annotate('', xy=(x + x_step, 0.5), xytext=(x + x_step * 0.85, 0.5),
                       arrowprops=dict(arrowstyle='->', lw=2, color='black'))

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    if title:
        ax.set_title(title, fontsize=12, fontweight='bold')

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_architecture_diagram(title, components):
    """Create architecture diagram with stacked components."""
    fig, ax = plt.subplots(figsize=(8, 6))

    colors = ['#3498db', '#e74c3c', '#2ecc71', '#f39c12', '#9b59b6']
    n = len(components)
    y_step = 0.7 / n
    y_start = 0.15

    for i, comp in enumerate(components):
        y = y_start + i * y_step
        color = colors[i % len(colors)]
        rect = patches.FancyBboxPatch((0.2, y), 0.6, y_step * 0.8,
                                     boxstyle="round,pad=0.01",
                                     edgecolor='black', facecolor=color,
                                     alpha=0.4, linewidth=2)
        ax.add_patch(rect)
        ax.text(0.5, y + y_step * 0.4, comp, ha='center', va='center',
               fontsize=10, fontweight='bold')

        if i < n - 1:
            ax.annotate('', xy=(0.5, y_start + (i+1) * y_step),
                       xytext=(0.5, y + y_step * 0.8),
                       arrowprops=dict(arrowstyle='->', lw=2, color='black'))

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    ax.set_title(title, fontsize=12, fontweight='bold')

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_comparison_chart(models, scores, title):
    """Create horizontal bar chart."""
    fig, ax = plt.subplots(figsize=(9, 5))

    colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DFE6E9', '#6C5CE7']
    bars = ax.barh(models, scores, color=colors[:len(models)], alpha=0.7,
                   edgecolor='black', linewidth=1.2)

    for i, (bar, val) in enumerate(zip(bars, scores)):
        ax.text(val + 1, i, f'{val:.1f}', va='center', fontweight='bold', fontsize=9)

    ax.set_xlabel('Performance', fontsize=10, fontweight='bold')
    ax.set_title(title, fontsize=12, fontweight='bold')
    ax.grid(axis='x', alpha=0.3, linestyle='--')

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_attention_visual():
    """Create attention heatmap."""
    fig, ax = plt.subplots(figsize=(7, 5))

    size = 8
    attention = np.random.rand(size, size)
    for i in range(size):
        attention[i, i] += 0.5
        if i > 0:
            attention[i, i-1] += 0.3
        if i < size - 1:
            attention[i, i+1] += 0.3
    attention = attention / attention.max()

    im = ax.imshow(attention, cmap='YlOrRd', aspect='auto')
    ax.set_title('Cross-Attention Pattern', fontsize=12, fontweight='bold')
    ax.set_xlabel('Keys (Image)', fontsize=9)
    ax.set_ylabel('Queries (Text)', fontsize=9)

    plt.colorbar(im, ax=ax, label='Attention Weight')

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_fusion_diagram():
    """Create fusion strategies diagram."""
    fig, axes = plt.subplots(1, 3, figsize=(11, 3.5))
    strategies = ['Early Fusion', 'Late Fusion', 'Hybrid Fusion']

    for idx, (ax, strategy) in enumerate(zip(axes, strategies)):
        if idx == 0:
            ax.add_patch(patches.Rectangle((0.2, 0.7), 0.25, 0.15, facecolor='#FF6B6B', alpha=0.5))
            ax.text(0.325, 0.775, 'Vision', ha='center', fontsize=8)
            ax.add_patch(patches.Rectangle((0.55, 0.7), 0.25, 0.15, facecolor='#4ECDC4', alpha=0.5))
            ax.text(0.675, 0.775, 'Text', ha='center', fontsize=8)
            ax.annotate('', xy=(0.5, 0.45), xytext=(0.325, 0.7),
                       arrowprops=dict(arrowstyle='->', lw=1.5))
            ax.annotate('', xy=(0.5, 0.45), xytext=(0.675, 0.7),
                       arrowprops=dict(arrowstyle='->', lw=1.5))
            ax.add_patch(patches.Rectangle((0.35, 0.3), 0.3, 0.2,
                                          facecolor='#9B59B6', alpha=0.5))
            ax.text(0.5, 0.4, 'Joint\nProcess', ha='center', fontsize=8)
        elif idx == 1:
            ax.add_patch(patches.Rectangle((0.1, 0.5), 0.3, 0.3,
                                          facecolor='#FF6B6B', alpha=0.3))
            ax.text(0.25, 0.7, 'Vision', ha='center', fontsize=8)
            ax.add_patch(patches.Rectangle((0.6, 0.5), 0.3, 0.3,
                                          facecolor='#4ECDC4', alpha=0.3))
            ax.text(0.75, 0.7, 'Text', ha='center', fontsize=8)
            ax.annotate('', xy=(0.5, 0.25), xytext=(0.25, 0.5),
                       arrowprops=dict(arrowstyle='->', lw=1.5))
            ax.annotate('', xy=(0.5, 0.25), xytext=(0.75, 0.5),
                       arrowprops=dict(arrowstyle='->', lw=1.5))
        else:
            ax.add_patch(patches.Rectangle((0.2, 0.7), 0.25, 0.12, facecolor='#FF6B6B', alpha=0.5))
            ax.text(0.325, 0.76, 'Vision', ha='center', fontsize=8)
            ax.add_patch(patches.Rectangle((0.55, 0.7), 0.25, 0.12, facecolor='#4ECDC4', alpha=0.5))
            ax.text(0.675, 0.76, 'Text', ha='center', fontsize=8)
            ax.add_patch(patches.Rectangle((0.3, 0.45), 0.4, 0.1,
                                          facecolor='#FFA07A', alpha=0.5))
            ax.text(0.5, 0.5, 'Fusion', ha='center', fontsize=8)

        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        ax.set_title(strategy, fontsize=9, fontweight='bold')

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_lora_diagram():
    """Create LoRA diagram."""
    fig, ax = plt.subplots(figsize=(9, 5))

    ax.add_patch(patches.Rectangle((0.1, 0.4), 0.2, 0.25,
                                   facecolor='#3498db', alpha=0.5, linewidth=2, edgecolor='black'))
    ax.text(0.2, 0.525, 'W\n(Frozen)', ha='center', va='center', fontsize=11, fontweight='bold')
    ax.text(0.35, 0.525, '+', ha='center', va='center', fontsize=18, fontweight='bold')

    ax.add_patch(patches.Rectangle((0.45, 0.45), 0.08, 0.22,
                                   facecolor='#e74c3c', alpha=0.5, linewidth=2, edgecolor='black'))
    ax.text(0.49, 0.56, 'B', ha='center', va='center', fontsize=11, fontweight='bold')

    ax.text(0.57, 0.56, '×', ha='center', va='center', fontsize=16, fontweight='bold')

    ax.add_patch(patches.Rectangle((0.62, 0.5), 0.12, 0.08,
                                   facecolor='#2ecc71', alpha=0.5, linewidth=2, edgecolor='black'))
    ax.text(0.68, 0.54, 'A', ha='center', va='center', fontsize=11, fontweight='bold')

    ax.text(0.5, 0.25, "W' = W + B × A", ha='center', fontsize=12, fontweight='bold',
           bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))

    ax.text(0.5, 0.12, "✓ Train only 1-2% of parameters", ha='center', fontsize=9)
    ax.text(0.5, 0.05, "✓ Fast training, low memory", ha='center', fontsize=9)

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 0.75)
    ax.axis('off')
    ax.set_title('LoRA: Low-Rank Adaptation', fontsize=12, fontweight='bold')

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_clip_architecture():
    """CLIP architecture diagram."""
    fig, ax = plt.subplots(figsize=(9, 5))

    ax.add_patch(patches.FancyBboxPatch((0.05, 0.5), 0.15, 0.3, boxstyle="round",
                                       edgecolor='#FF6B6B', facecolor='#FF6B6B', alpha=0.3, linewidth=2))
    ax.text(0.125, 0.65, 'Image\nEncoder', ha='center', va='center', fontsize=9, fontweight='bold')

    ax.add_patch(patches.FancyBboxPatch((0.05, 0.1), 0.15, 0.3, boxstyle="round",
                                       edgecolor='#4ECDC4', facecolor='#4ECDC4', alpha=0.3, linewidth=2))
    ax.text(0.125, 0.25, 'Text\nEncoder', ha='center', va='center', fontsize=9, fontweight='bold')

    ax.add_patch(patches.FancyBboxPatch((0.3, 0.5), 0.15, 0.3, boxstyle="round",
                                       edgecolor='#9B59B6', facecolor='#9B59B6', alpha=0.3, linewidth=2))
    ax.text(0.375, 0.65, 'Image\nProj', ha='center', va='center', fontsize=9, fontweight='bold')

    ax.add_patch(patches.FancyBboxPatch((0.3, 0.1), 0.15, 0.3, boxstyle="round",
                                       edgecolor='#9B59B6', facecolor='#9B59B6', alpha=0.3, linewidth=2))
    ax.text(0.375, 0.25, 'Text\nProj', ha='center', va='center', fontsize=9, fontweight='bold')

    ax.add_patch(patches.FancyBboxPatch((0.6, 0.25), 0.3, 0.4, boxstyle="round",
                                       edgecolor='#2ECC71', facecolor='#2ECC71', alpha=0.3, linewidth=3))
    ax.text(0.75, 0.5, 'Joint\nEmbedding\nSpace', ha='center', va='center', fontsize=10, fontweight='bold')
    ax.text(0.75, 0.35, 'Contrastive\nLearning', ha='center', va='center', fontsize=8, style='italic')

    for y_start, y_end in [(0.65, 0.5), (0.25, 0.35)]:
        ax.annotate('', xy=(0.3, y_start), xytext=(0.2, y_start),
                   arrowprops=dict(arrowstyle='->', lw=1.5))
        ax.annotate('', xy=(0.6, y_end), xytext=(0.45, y_start),
                   arrowprops=dict(arrowstyle='->', lw=1.5))

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 0.9)
    ax.axis('off')
    ax.set_title('CLIP Architecture', fontsize=12, fontweight='bold')

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
    buf.seek(0)
    plt.close()
    return buf

def create_presentation():
    """Create complete presentation with images."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    CLEMSON_ORANGE = RGBColor(246, 103, 51)
    CLEMSON_PURPLE = RGBColor(82, 45, 128)
    DARK_GRAY = RGBColor(51, 51, 51)

    def add_image_to_slide(slide, image_stream, left, top, width):
        """Add image to slide."""
        slide.shapes.add_picture(image_stream, left, top, width)

    def add_title_slide(title, subtitle="", image_stream=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = CLEMSON_ORANGE

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(22)
            subtitle_para.font.color.rgb = DARK_GRAY

        if image_stream:
            add_image_to_slide(slide, image_stream, Inches(3), Inches(4.8), Inches(4))

        return slide

    def add_content_slide(title, content_items, image_stream=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(26)
        title_para.font.bold = True
        title_para.font.color.rgb = CLEMSON_PURPLE

        if image_stream:
            content_left, content_width = Inches(0.5), Inches(4.8)
            img_left, img_width = Inches(5.5), Inches(4.2)
        else:
            content_left, content_width = Inches(0.7), Inches(8.6)

        content_box = slide.shapes.add_textbox(content_left, Inches(1.1), content_width, Inches(6))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]

            if isinstance(item, tuple):
                p.text = item[0]
                p.level = item[1]
            else:
                p.text = item
                p.level = 0

            p.font.size = Pt(15) if p.level == 0 else Pt(13)
            p.space_before = Pt(5)
            p.font.color.rgb = DARK_GRAY

        if image_stream:
            add_image_to_slide(slide, image_stream, img_left, Inches(1.5), img_width)

        return slide

    def add_section_slide(section_title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = CLEMSON_PURPLE

        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        return slide

    print("\n" + "="*70)
    print("GENERATING MULTIMODAL LLM PRESENTATION WITH IMAGES")
    print("="*70 + "\n")

    # Slide 1
    log_status("Title slide - Adding multimodal learning diagram")
    img = create_diagram_multimodal_learning()
    add_title_slide("Multimodal Large Language Models",
                    "Applied Data Science - Clemson University", img)

    # Slide 2
    log_status("Course Context - Adding course flow diagram")
    img = create_simple_diagram("Course Flow", ["Transformers", "LLMs", "Multimodal\nLLMs", "Applications"])
    add_content_slide("Course Context", [
        "Part of Applied Data Science course",
        "Master's program in Computer Science",
        "Prerequisites: Understanding of Transformer architecture",
        ("Attention mechanisms", 1),
        ("Self-attention and cross-attention", 1),
        "Building on previous lessons on LLMs"
    ], img)

    # Slide 3
    log_status("Learning Objectives - Adding objectives checklist visual")
    img = create_simple_diagram("Learning Path", ["Theory", "Architectures", "Training", "Practice"])
    add_content_slide("Learning Objectives", [
        "Understand multimodal learning fundamentals",
        "Master single-modality representations",
        "Learn multimodal fusion strategies",
        "Study state-of-the-art architectures",
        "Understand training vs. fine-tuning",
        "Apply on Palmetto cluster"
    ], img)

    # Slide 4
    log_status("What is Multimodal Learning - Adding concept diagram")
    img = create_diagram_multimodal_learning()
    add_content_slide("What is Multimodal Learning?", [
        "Learning from multiple data modalities simultaneously",
        "Modalities: Different types of data sources",
        ("Vision (images, videos)", 1),
        ("Text (documents, captions)", 1),
        ("Audio (speech, sound)", 1),
        "Goal: Build unified representations",
        "Enable cross-modal understanding"
    ], img)

    # Slide 5
    log_status("Why Multimodal Models - Adding applications diagram")
    img = create_simple_diagram("Applications", ["VQA", "Captioning", "Generation", "Robotics"])
    add_content_slide("Why Multimodal Models?", [
        "Human perception is inherently multimodal",
        "Richer understanding through multiple perspectives",
        "Key applications:",
        ("Visual Question Answering", 1),
        ("Image/Video captioning", 1),
        ("Text-to-image generation", 1),
        ("Embodied AI and robotics", 1),
        "Better generalization and robustness"
    ], img)

    # Slide 6
    log_status("Historical Context - Adding timeline visualization")
    img = create_simple_diagram("Evolution", ["2017\nAttention", "2019\nBERT+Vision", "2021\nCLIP", "2023\nGPT-4V"])
    add_content_slide("Historical Context", [
        "2010s: Separate models per modality",
        "2017: Attention mechanisms enable fusion",
        "2019-2020: BERT + Vision (ViLBERT, LXMERT)",
        "2021: CLIP revolutionizes vision-language",
        "2022-2023: Large-scale models (Flamingo, GPT-4V)",
        "2023-2024: Open-source alternatives (LLaVA, BLIP-2)",
        "2024+: Any-to-any multimodal models"
    ], img)

    # Slide 7
    log_status("Challenges - Adding challenges diagram")
    img = create_simple_diagram("Key Challenges", ["Modality\nGap", "Alignment", "Fusion", "Scale"])
    add_content_slide("Challenges in Multimodal Learning", [
        "Modality gap: Different data distributions",
        "Representation learning: Finding common space",
        "Alignment: Matching corresponding elements",
        "Fusion strategies: How to combine?",
        "Data requirements: Need paired data",
        "Computational cost: Processing multiple modalities",
        "Evaluation: Measuring performance"
    ], img)

    # Slide 8 - Section
    log_status("Section: Single Modality Foundations")
    add_section_slide("Part I: Single Modality Foundations")

    # Slide 9
    log_status("Vision Modality - Adding vision pipeline diagram")
    img = create_architecture_diagram("Vision Pipeline", ["Image\nInput", "Patch\nEmbedding", "Transformer\nEncoder", "Features"])
    add_content_slide("Single Modality: Vision", [
        "Image representation: Pixels in grids",
        "Convolutional Neural Networks (CNNs):",
        ("Learn spatial hierarchies", 1),
        ("Translation invariance", 1),
        "Vision Transformers (ViT):",
        ("Split image into patches", 1),
        ("Apply transformer to patches", 1),
        ("Better scalability with data", 1),
        "Output: High-dimensional feature vectors"
    ], img)

    # Slide 10
    log_status("Vision Transformers - Adding ViT architecture")
    img = create_architecture_diagram("ViT Architecture", ["Image\nPatches", "Linear\nProjection", "Positional\nEmbedding", "Transformer\nLayers", "[CLS]\nToken"])
    add_content_slide("Vision Transformers (ViT)", [
        "Pure transformer applied to images",
        "Process:",
        ("1. Split image into 16×16 patches", 1),
        ("2. Flatten and linearly embed", 1),
        ("3. Add positional embeddings", 1),
        ("4. Feed through transformer", 1),
        ("5. Extract [CLS] token", 1),
        "Advantages: Scalable, less bias",
        "Used in CLIP, BLIP, modern models"
    ], img)

    # Slide 11
    log_status("Text Modality - Adding text processing pipeline")
    img = create_architecture_diagram("Text Pipeline", ["Text\nInput", "Tokenization", "Embedding", "Transformer", "Contextualized\nOutput"])
    add_content_slide("Single Modality: Text", [
        "Tokenization: Convert text to tokens",
        "Embedding: Map tokens to vectors",
        "Transformers for text:",
        ("BERT: Bidirectional encoder", 1),
        ("GPT: Autoregressive decoder", 1),
        ("T5: Encoder-decoder", 1),
        "Contextual representations",
        "Pre-training on large text corpora",
        "Output: Sequence of embeddings"
    ], img)

    # Slide 12
    log_status("Audio Modality - Adding audio processing diagram")
    img = create_architecture_diagram("Audio Pipeline", ["Waveform", "Spectrogram", "Features", "Model", "Embeddings"])
    add_content_slide("Single Modality: Audio", [
        "Raw audio: Waveform representation",
        "Feature extraction:",
        ("Mel-spectrograms: Time-frequency", 1),
        ("MFCCs: Cepstral coefficients", 1),
        "Models:",
        ("Wav2Vec 2.0: Self-supervised speech", 1),
        ("Whisper: Robust speech recognition", 1),
        ("AudioMAE: Masked autoencoding", 1),
        "Output: Audio embeddings"
    ], img)

    # Slide 13
    log_status("Video Modality - Adding video processing diagram")
    img = create_architecture_diagram("Video Pipeline", ["Frames", "Spatial\nEncoding", "Temporal\nEncoding", "Fusion", "Video\nFeatures"])
    add_content_slide("Single Modality: Video", [
        "Video = Images + temporal information",
        "Spatial encoding: Per-frame CNN/ViT",
        "Temporal encoding:",
        ("3D convolutions", 1),
        ("Recurrent networks (LSTM)", 1),
        ("Temporal transformers", 1),
        "Challenges: Computational cost",
        "Models: TimeSformer, VideoMAE",
        "Output: Spatio-temporal features"
    ], img)

    # Slide 14
    log_status("Modality Representations - Adding representation space diagram")
    img = create_diagram_multimodal_learning()
    add_content_slide("Modality Representations", [
        "Each modality has unique characteristics",
        "Vision: Spatial structure, continuous",
        "Text: Sequential, discrete tokens",
        "Audio: Temporal, continuous waveforms",
        "Goal: Map to common semantic space",
        "Representation learning objectives:",
        ("Reconstruction", 1),
        ("Contrastive learning", 1),
        ("Masked modeling", 1)
    ], img)

    # Slide 15
    log_status("Feature Extraction - Adding extraction pipeline")
    img = create_simple_diagram("Extraction Pipeline", ["Input", "Preprocess", "Encoder", "Projection", "Embeddings"])
    add_content_slide("Feature Extraction Pipeline", [
        "Input: Raw data in original modality",
        "Preprocessing:",
        ("Images: Resize, normalize", 1),
        ("Text: Tokenize, add special tokens", 1),
        ("Audio: Resample, spectrogram", 1),
        "Encoder: Modality-specific network",
        "Projection: Map to common dimension",
        "Normalization: L2 normalize for similarity"
    ], img)

    # Slide 16
    log_status("Embedding Spaces - Adding embedding space visualization")
    img = create_attention_visual()
    add_content_slide("Embedding Spaces", [
        "Embedding space: High-dimensional vectors",
        "Desirable properties:",
        ("Semantic similarity → Vector proximity", 1),
        ("Modality invariance", 1),
        ("Structured relationships", 1),
        "Joint embedding space:",
        ("Multiple modalities → same space", 1),
        ("Enable cross-modal retrieval", 1),
        "Learned through contrastive objectives"
    ], img)

    # Slide 17
    log_status("Alignment Problem - Adding alignment diagram")
    img = create_simple_diagram("Alignment Types", ["Global\nAlignment", "Local\nAlignment", "Temporal\nAlignment"])
    add_content_slide("The Alignment Problem", [
        "Core challenge: Align different modalities",
        "Types of alignment:",
        ("Global: Image-caption pairs", 1),
        ("Local: Objects-words", 1),
        ("Temporal: Video-transcript", 1),
        "Solutions:",
        ("Contrastive learning (CLIP)", 1),
        ("Cross-attention (Flamingo)", 1),
        ("Q-Former (BLIP-2)", 1),
        "Goal: Bridge the modality gap"
    ], img)

    # Slide 18 - Section
    log_status("Section: Multimodal Architectures")
    add_section_slide("Part II: Multimodal Architectures")

    # Slide 19
    log_status("Architecture Overview - Adding architecture taxonomy")
    img = create_simple_diagram("Architecture Types", ["Dual\nEncoder", "Fusion\nEncoder", "Encoder\nDecoder"])
    add_content_slide("Overview of Multimodal Architectures", [
        "Three main paradigms:",
        "1. Dual-encoder (CLIP, ALIGN)",
        ("Separate encoders per modality", 1),
        ("Joint embedding space", 1),
        "2. Fusion encoder (LXMERT, ViLBERT)",
        ("Cross-modal attention layers", 1),
        "3. Encoder-decoder (BLIP, Flamingo, LLaVA)",
        ("Visual encoder + Language model", 1),
        ("Generation capabilities", 1)
    ], img)

    # Slide 20
    log_status("Early Fusion - Adding early fusion diagram")
    img = create_fusion_diagram()
    add_content_slide("Early Fusion", [
        "Combine raw features at input level",
        "Concatenate or sum modality features",
        "Process jointly through unified network",
        "Advantages:",
        ("Maximum interaction", 1),
        ("Simple architecture", 1),
        "Disadvantages:",
        ("Hard to leverage pre-trained models", 1),
        ("Less flexible", 1)
    ], img)

    # Slide 21
    log_status("Late Fusion - Part of fusion strategies diagram shown")
    img = create_fusion_diagram()
    add_content_slide("Late Fusion", [
        "Process each modality independently",
        "Combine at decision/output level",
        "Examples: Ensemble approaches",
        "Advantages:",
        ("Easy to use pre-trained models", 1),
        ("Modular design", 1),
        ("Train modalities separately", 1),
        "Disadvantages:",
        ("Limited cross-modal interaction", 1),
        ("May miss complex relationships", 1)
    ], img)

    # Slide 22
    log_status("Hybrid Fusion - Complete fusion comparison")
    img = create_fusion_diagram()
    add_content_slide("Hybrid Fusion", [
        "Combine early and late fusion benefits",
        "Multi-stage interaction",
        "Example: Cross-attention at multiple layers",
        "Advantages:",
        ("Flexible architecture design", 1),
        ("Leverage pre-trained components", 1),
        ("Rich cross-modal interaction", 1),
        "Most modern architectures use hybrid",
        "Examples: BLIP-2, Flamingo, LLaVA"
    ], img)

    # Slide 23
    log_status("Cross-Attention - Adding attention visualization")
    img = create_attention_visual()
    add_content_slide("Cross-Attention Mechanisms", [
        "Key technique for multimodal fusion",
        "Standard attention: Q, K, V from same sequence",
        "Cross-attention: Q from one modality, K,V from another",
        "Example: Text queries vision features",
        ("Q = text embeddings", 1),
        ("K, V = image patch features", 1),
        "Bidirectional cross-attention for deep fusion",
        "Used in: Flamingo, BLIP, Q-Former"
    ], img)

    # Slide 24
    log_status("CLIP Architecture - Adding CLIP diagram")
    img = create_clip_architecture()
    add_content_slide("CLIP: Contrastive Language-Image Pre-training", [
        "Dual-encoder architecture (OpenAI, 2021)",
        "Components:",
        ("Image encoder: ViT or ResNet", 1),
        ("Text encoder: Transformer", 1),
        ("Project to same dimension", 1),
        "Training: Contrastive learning on 400M pairs",
        "Zero-shot capabilities",
        "Applications: Classification, retrieval"
    ], img)

    # Slide 25
    log_status("CLIP Training - Adding contrastive loss diagram")
    img = create_simple_diagram("Contrastive Learning", ["Positive\nPairs", "Maximize\nSimilarity", "Negative\nPairs", "Minimize\nSimilarity"])
    add_content_slide("CLIP Training Objective", [
        "Contrastive loss on batch of N pairs",
        "Maximize similarity of correct (image, text) pairs",
        "Minimize similarity of incorrect pairs",
        "Loss: InfoNCE",
        "Temperature scaling: τ learned parameter",
        "Symmetric: Both image→text and text→image",
        "Result: Aligned vision-language space"
    ], img)

    # Slide 26
    log_status("ALIGN Architecture - Adding ALIGN comparison")
    img = create_comparison_chart(["CLIP", "ALIGN"], [68.7, 76.4], "Zero-Shot Performance")
    add_content_slide("ALIGN: A Large-scale ImaGe and Noisy text", [
        "Similar to CLIP (Google, 2021)",
        "Key difference: Noisy web data (1.8B pairs)",
        "No careful curation",
        "Components:",
        ("EfficientNet for images", 1),
        ("BERT for text", 1),
        "Same contrastive objective",
        "Insight: Scale > data quality",
        "Comparable or better results"
    ], img)

    # Slide 27
    log_status("BLIP Architecture - Adding BLIP diagram")
    img = create_architecture_diagram("BLIP", ["Image\nEncoder", "Text\nEncoder", "Text\nDecoder", "Multi-task\nLearning"])
    add_content_slide("BLIP: Bootstrapping Language-Image Pre-training", [
        "Encoder-decoder architecture (Salesforce, 2022)",
        "Three components:",
        ("1. Image encoder: ViT", 1),
        ("2. Text encoder: BERT-like", 1),
        ("3. Text decoder: GPT-like", 1),
        "Multi-task learning:",
        ("Contrastive learning", 1),
        ("Image-text matching", 1),
        ("Image-conditioned LM", 1),
        "CapFilt: Synthetic captions"
    ], img)

    # Slide 28
    log_status("BLIP-2 Q-Former - Adding Q-Former architecture")
    img = create_architecture_diagram("BLIP-2", ["Frozen\nViT", "Q-Former\n(Learnable)", "Frozen\nLLM", "Efficient!"])
    add_content_slide("BLIP-2: Querying Transformer", [
        "Efficient architecture (Salesforce, 2023)",
        "Key innovation: Q-Former",
        ("Lightweight module between frozen encoders", 1),
        ("Learnable query tokens", 1),
        ("Cross-attention to vision", 1),
        "Two-stage training:",
        ("Stage 1: Vision-language learning", 1),
        ("Stage 2: Vision-to-language generation", 1),
        "Much more parameter-efficient"
    ], img)

    # Slide 29
    log_status("Flamingo Architecture - Adding Flamingo diagram")
    img = create_architecture_diagram("Flamingo", ["Vision\nEncoder", "Perceiver\nResampler", "Gated\nXATTN", "LLM\n(70B)"])
    add_content_slide("Flamingo: Visual Language Model", [
        "Few-shot learning VLM (DeepMind, 2022)",
        "Architecture:",
        ("Vision encoder: Frozen pre-trained", 1),
        ("Perceiver resampler: Compress features", 1),
        ("LLM: Frozen Chinchilla 70B", 1),
        ("Cross-attention: Gated XATTN-DENSE", 1),
        "Interleaved image-text input",
        "In-context learning: Few-shot prompting",
        "Powerful few-shot visual understanding"
    ], img)

    # Slide 30
    log_status("GPT-4V Overview - Adding GPT-4V capabilities")
    img = create_simple_diagram("GPT-4V Capabilities", ["VQA", "OCR", "Reasoning", "Code\nGen"])
    add_content_slide("GPT-4V: Multimodal GPT-4", [
        "Proprietary multimodal model (OpenAI, 2023)",
        "Accepts images and text as input",
        "Architecture details not disclosed",
        "Likely components:",
        ("Enhanced CLIP-like vision encoder", 1),
        ("Large language model: GPT-4", 1),
        ("Sophisticated fusion", 1),
        "Capabilities: VQA, OCR, reasoning, code",
        "State-of-the-art performance"
    ], img)

    # Slide 31
    log_status("LLaVA Architecture - Adding LLaVA diagram")
    img = create_architecture_diagram("LLaVA", ["CLIP\nViT", "Projection\nLayer", "Vicuna\nLLM", "Simple &\nEffective"])
    add_content_slide("LLaVA: Large Language and Vision Assistant", [
        "Open-source VLM (Liu et al., 2023)",
        "Simple architecture:",
        ("Vision encoder: CLIP ViT-L/14", 1),
        ("Projection: Linear or MLP", 1),
        ("LLM: Vicuna (LLaMA-based)", 1),
        "Training data: GPT-4 generated instructions",
        "Two-stage training:",
        ("Pre-training: Projection only", 1),
        ("Fine-tuning: Full model or LoRA", 1),
        "Excellent performance, modest compute"
    ], img)

    # Slide 32
    log_status("LLaVA 1.5 - Adding improvements comparison")
    img = create_comparison_chart(["LLaVA 1.0", "LLaVA 1.5", "GPT-4V"], [70, 85, 90], "Benchmark Performance")
    add_content_slide("LLaVA 1.5 Improvements", [
        "Enhanced version (2023)",
        "Improvements:",
        ("MLP projection: 2-layer vs 1-layer", 1),
        ("Higher resolution: 336×336 vs 224", 1),
        ("Better instruction data", 1),
        ("Academic task data: VQA, OCR", 1),
        "Unified training recipe",
        "Near GPT-4V on some tasks",
        "Efficient: Fine-tune on single GPU"
    ], img)

    # Slide 33
    log_status("Gemini Architecture - Adding Gemini overview")
    img = create_architecture_diagram("Gemini", ["Native\nMultimodal", "All\nModalities", "Long\nContext", "State-of-\nthe-Art"])
    add_content_slide("Gemini: Google's Multimodal Model", [
        "Native multimodal model (Google, 2023)",
        "Three sizes: Nano, Pro, Ultra",
        "Key: Natively multimodal from scratch",
        ("Not adapting separate models", 1),
        ("Joint training on all modalities", 1),
        "Capabilities:",
        ("Image, video, audio, text", 1),
        ("Long context: Up to 1M tokens", 1),
        ("Complex reasoning", 1),
        "State-of-the-art benchmarks"
    ], img)

    # Slide 34
    log_status("Any-to-Any Models - Adding any-to-any diagram")
    img = create_diagram_multimodal_learning()
    add_content_slide("Any-to-Any Multimodal Models", [
        "Next generation: Unified for all modalities",
        "Examples:",
        ("NExT-GPT: Any-to-any generation", 1),
        ("CoDi: Composable diffusion", 1),
        ("Unified-IO 2: Single model, any task", 1),
        "Capabilities:",
        ("Text → Image, Audio, Video", 1),
        ("Image → Text, Audio", 1),
        ("Any modality → Any modality", 1),
        "Challenges: Training complexity, data"
    ], img)

    # Slide 35
    log_status("Unified Models - Adding unified architecture")
    img = create_architecture_diagram("Unified Model", ["Tokenize\nAll", "Shared\nEncoder", "Shared\nDecoder", "Any→Any"])
    add_content_slide("Unified Multimodal Models", [
        "Goal: Single model for all modalities",
        "Approaches:",
        ("Tokenize all modalities (Unified-IO)", 1),
        ("Shared encoder-decoder", 1),
        ("Modality-specific adapters", 1),
        "Benefits:",
        ("Simplified deployment", 1),
        ("Cross-modal transfer", 1),
        ("Efficient inference", 1),
        "Research direction: Foundation models"
    ], img)

    # Slide 36
    log_status("Architecture Comparison - Adding performance chart")
    img = create_comparison_chart(
        ["CLIP", "BLIP", "BLIP-2", "Flamingo", "LLaVA", "GPT-4V"],
        [68.7, 82.3, 84.5, 85.8, 88.5, 92.1],
        "Model Performance Comparison"
    )
    add_content_slide("Architecture Comparison", [
        "CLIP: Fast, dual-encoder, zero-shot",
        "BLIP: Versatile, generation + understanding",
        "BLIP-2: Efficient, frozen components",
        "Flamingo: Few-shot, interleaved input",
        "LLaVA: Open-source, instruction-tuned",
        "GPT-4V: State-of-the-art, proprietary",
        "Selection: Task, compute, data, performance"
    ], img)

    # Slide 37
    log_status("Design Principles - Adding principles checklist")
    img = create_simple_diagram("Design Principles", ["Leverage\nPre-trained", "Freeze\nLarge", "Efficient\nFusion", "Multi-task"])
    add_content_slide("Design Principles", [
        "1. Leverage pre-trained models when possible",
        "2. Freeze large components to reduce compute",
        "3. Use efficient fusion (Q-Former, projectors)",
        "4. Multi-task learning for generalization",
        "5. High-quality instruction data crucial",
        "6. Balance model size vs. performance",
        "7. Consider inference efficiency",
        "8. Open-source enables research"
    ], img)

    # Slide 38 - Section
    log_status("Section: Training Strategies")
    add_section_slide("Part III: Training Strategies")

    # Slide 39
    log_status("Training from Scratch - Adding training considerations")
    img = create_simple_diagram("Training Considerations", ["Data\nScale", "Compute\nCost", "Time\nRequired", "Expertise"])
    add_content_slide("Training from Scratch", [
        "When to train from scratch?",
        ("Novel architectures", 1),
        ("Unique data distribution", 1),
        ("Proprietary requirements", 1),
        "Challenges:",
        ("Massive data (millions of pairs)", 1),
        ("Extensive compute (1000s GPU hours)", 1),
        ("Careful hyperparameter tuning", 1),
        "Alternative: Fine-tune pre-trained models"
    ], img)

    # Slide 40
    log_status("Data Requirements - Adding data scale visualization")
    img = create_comparison_chart(
        ["COCO", "CC3M", "LAION-400M", "LAION-5B"],
        [0.3, 3, 400, 5000],
        "Dataset Sizes (Millions)"
    )
    add_content_slide("Data Requirements", [
        "Scale: Millions/billions of pairs",
        "Types of data:",
        ("Image-caption pairs", 1),
        ("Interleaved image-text", 1),
        ("Video-subtitle pairs", 1),
        ("Question-answer pairs", 1),
        "Data quality vs. quantity:",
        ("CLIP: Curated 400M", 1),
        ("ALIGN: Noisy 1.8B", 1),
        "Diversity crucial for generalization"
    ], img)

    # Slide 41
    log_status("Dataset Examples - Adding dataset overview")
    img = create_simple_diagram("Dataset Types", ["Image-Text", "VQA", "Video", "Dense\nAnnotations"])
    add_content_slide("Key Datasets", [
        "Image-Text:",
        ("COCO: 330K images with captions", 1),
        ("Visual Genome: 108K dense annotations", 1),
        ("Conceptual Captions: 3.3M pairs", 1),
        ("LAION-400M/5B: Large-scale web", 1),
        "VQA:",
        ("VQAv2: 1.1M questions", 1),
        ("GQA: 22M questions with graphs", 1),
        "Video: WebVid, HowTo100M, YT-Temporal"
    ], img)

    # Slide 42
    log_status("Pre-training Objectives - Adding objectives diagram")
    img = create_simple_diagram("Training Objectives", ["Contrastive", "Masked\nModeling", "ITM", "Generative"])
    add_content_slide("Pre-training Objectives", [
        "1. Contrastive Learning:",
        ("Align positive pairs, separate negatives", 1),
        ("InfoNCE loss", 1),
        "2. Masked Language/Image Modeling:",
        ("Mask tokens/patches, predict them", 1),
        ("Learn contextualized representations", 1),
        "3. Image-Text Matching:",
        ("Binary: match or not", 1),
        "4. Generative objectives:",
        ("Caption/image generation", 1)
    ], img)

    # Slide 43
    log_status("Contrastive Learning - Adding contrastive learning visual")
    img = create_simple_diagram("Contrastive Loss", ["Positive\nPairs", "Pull\nTogether", "Negative\nPairs", "Push\nApart"])
    add_content_slide("Contrastive Learning", [
        "Core idea: Similar items close, dissimilar far",
        "For multimodal: Match corresponding modalities",
        "Positive pairs: (image, caption) from same example",
        "Negative pairs: All other combinations",
        "Loss function: InfoNCE",
        ("Numerator: Similarity of positive", 1),
        ("Denominator: Sum of all similarities", 1),
        "Temperature scaling: Control sharpness",
        "Effective for aligned embeddings"
    ], img)

    # Slide 44
    log_status("Masked Modeling - Adding masked modeling diagram")
    img = create_architecture_diagram("Masked Modeling", ["Input", "Mask\nTokens", "Encode", "Predict\nMasked", "Loss"])
    add_content_slide("Masked Language/Image Modeling", [
        "Masked Language Modeling (MLM):",
        ("Mask random tokens in text", 1),
        ("Predict from context", 1),
        ("Used in BERT, VilBERT", 1),
        "Masked Image Modeling (MIM):",
        ("Mask random patches", 1),
        ("Reconstruct patches", 1),
        ("Used in MAE, SimMIM", 1),
        "Multimodal: Mask in both",
        "Learn robust cross-modal representations"
    ], img)

    # Slide 45
    log_status("Image-Text Matching - Adding ITM diagram")
    img = create_simple_diagram("ITM Task", ["Image +\nText", "Cross\nAttention", "Match?\n(Yes/No)", "Binary\nLoss"])
    add_content_slide("Image-Text Matching (ITM)", [
        "Binary classification task",
        "Input: Image and text pair",
        "Output: Match (1) or not (0)",
        "Requires cross-modal interaction",
        ("Usually through cross-attention", 1),
        "Negative sampling strategies:",
        ("Random negatives from batch", 1),
        ("Hard negatives: Similar but wrong", 1),
        "Used alongside contrastive learning",
        "Improves fine-grained alignment"
    ], img)

    # Slide 46
    log_status("Training Stability - Adding stability tips")
    img = create_simple_diagram("Stability Tips", ["Gradient\nClipping", "Layer\nNorm", "Warmup", "Mixed\nPrecision"])
    add_content_slide("Training Stability", [
        "Challenges:",
        ("Large batch sizes needed", 1),
        ("Gradient explosion/vanishing", 1),
        ("Modality imbalance", 1),
        "Solutions:",
        ("Gradient clipping", 1),
        ("Layer normalization", 1),
        ("Warmup learning rate schedule", 1),
        ("Mixed precision (FP16/BF16)", 1),
        ("Careful initialization", 1),
        "Monitor: Loss curves, gradient norms"
    ], img)

    # Slide 47
    log_status("Computational Requirements - Adding compute comparison")
    img = create_comparison_chart(
        ["CLIP", "Flamingo", "LLaVA"],
        [592*12, 2000, 8*1],
        "GPU-Days Required"
    )
    add_content_slide("Computational Requirements", [
        "Training large models is expensive",
        "Examples:",
        ("CLIP: 592 V100 GPUs, 12 days", 1),
        ("Flamingo: TPUv4 pods, weeks", 1),
        ("LLaVA: 8 A100 GPUs, ~1 day", 1),
        "Factors: Model size, batch, data, resolution",
        "Cost reduction:",
        ("Freeze components", 1),
        ("Lower resolution in pre-training", 1),
        ("Use smaller models for research", 1)
    ], img)

    # Slide 48
    log_status("Distributed Training - Adding distributed training diagram")
    img = create_simple_diagram("Parallelism Types", ["Data\nParallel", "Model\nParallel", "Pipeline\nParallel", "Tensor\nParallel"])
    add_content_slide("Distributed Training", [
        "Essential for large-scale training",
        "Strategies:",
        ("Data parallelism: Split data across GPUs", 1),
        ("Model parallelism: Split model across GPUs", 1),
        ("Pipeline parallelism: Split layers", 1),
        ("Tensor parallelism: Split tensors", 1),
        "Frameworks:",
        ("PyTorch DDP, FSDP", 1),
        ("DeepSpeed, Megatron-LM", 1),
        "Communication overhead: Critical bottleneck"
    ], img)

    # Slide 49
    log_status("Training Pipeline - Adding pipeline diagram")
    img = create_simple_diagram("Pipeline", ["Data\nPrep", "Init\nModel", "Train\nLoop", "Validate", "Save"])
    add_content_slide("Training Pipeline", [
        "1. Data preparation:",
        ("Collect and clean data", 1),
        ("Create train/val/test splits", 1),
        "2. Preprocessing:",
        ("Tokenization, resizing", 1),
        ("Augmentation strategies", 1),
        "3. Model initialization:",
        ("Load pre-trained weights", 1),
        "4. Training loop:",
        ("Forward, loss, backward, update", 1),
        "5. Evaluation and checkpointing"
    ], img)

    # Slide 50
    log_status("Evaluation Metrics - Adding metrics overview")
    img = create_simple_diagram("Metrics", ["Retrieval\n(Recall@K)", "VQA\n(Accuracy)", "Captioning\n(CIDEr)", "Classification\n(Acc)"])
    add_content_slide("Evaluation Metrics", [
        "Image-Text Retrieval:",
        ("Recall@K, Mean Rank", 1),
        "VQA:",
        ("Accuracy, VQA score", 1),
        "Image Captioning:",
        ("BLEU, METEOR, CIDEr, SPICE", 1),
        "Zero-shot Classification:",
        ("Top-1/Top-5 accuracy", 1),
        "General:",
        ("Cross-modal similarity", 1),
        ("Human evaluation", 1)
    ], img)

    # Slide 51
    log_status("Benchmarks - Adding benchmark list")
    img = create_simple_diagram("Key Benchmarks", ["ImageNet", "VQAv2", "COCO", "MMBench"])
    add_content_slide("Key Benchmarks", [
        "Classification: ImageNet, CIFAR",
        "VQA: VQAv2, GQA, OK-VQA, VizWiz",
        "Captioning: COCO Captions, NoCaps",
        "Retrieval: COCO, Flickr30K",
        "Reasoning: NLVR2, CLEVR",
        "OCR: TextVQA, DocVQA, Infographics",
        "Video: MSRVTT, ActivityNet",
        "Comprehensive: MMBench, SEED-Bench, MM-Vet"
    ], img)

    # Slide 52 - Section
    log_status("Section: Fine-tuning Strategies")
    add_section_slide("Part IV: Fine-tuning Strategies")

    # Slide 53
    log_status("Fine-tuning Overview - Adding fine-tuning benefits")
    img = create_comparison_chart(
        ["Data Required", "Compute Cost", "Training Time", "Performance"],
        [10, 5, 8, 95],
        "Fine-tuning vs Training from Scratch"
    )
    add_content_slide("Fine-tuning Overview", [
        "Most practical approach",
        "Start with pre-trained multimodal model",
        "Adapt to specific task or domain",
        "Advantages:",
        ("Less data required", 1),
        ("Less compute required", 1),
        ("Faster convergence", 1),
        ("Better performance", 1),
        "Types: Full vs. parameter-efficient"
    ], img)

    # Slide 54
    log_status("Transfer Learning - Adding transfer learning diagram")
    img = create_simple_diagram("Transfer Learning", ["Pre-trained\nKnowledge", "Task-specific\nData", "Fine-tune", "Specialized\nModel"])
    add_content_slide("Transfer Learning", [
        "Leverage knowledge from pre-training",
        "Pre-trained model already understands:",
        ("Visual concepts", 1),
        ("Language semantics", 1),
        ("Cross-modal alignment", 1),
        "Fine-tuning teaches:",
        ("Task-specific patterns", 1),
        ("Domain-specific knowledge", 1),
        ("Output format requirements", 1),
        "Critical: Choose appropriate base model"
    ], img)

    # Slide 55
    log_status("PEFT Overview - Adding PEFT comparison")
    img = create_comparison_chart(
        ["Full Fine-tuning", "LoRA", "Adapters", "Prompt Tuning"],
        [100, 1.5, 2, 0.1],
        "Trainable Parameters (%)"
    )
    add_content_slide("Parameter-Efficient Fine-tuning (PEFT)", [
        "Goal: Fine-tune with minimal parameters",
        "Why PEFT?",
        ("Large models too expensive", 1),
        ("Reduced memory", 1),
        ("Faster training", 1),
        ("Easy deployment (store adapters)", 1),
        "Methods:",
        ("LoRA, Adapters, Prompt tuning", 1),
        "Often achieves comparable performance"
    ], img)

    # Slide 56
    log_status("LoRA - Adding LoRA diagram")
    img = create_lora_diagram()
    add_content_slide("LoRA: Low-Rank Adaptation", [
        "Learn low-rank updates to weights",
        "W_new = W_frozen + B × A",
        ("W_frozen: Original weights", 1),
        ("B, A: Low-rank matrices (rank r)", 1),
        "Only train B and A (1-2% of parameters)",
        "Apply to:",
        ("Attention projections (Q, K, V)", 1),
        ("Feed-forward layers", 1),
        "Works well for LLM in VLMs",
        "Implementation: PEFT library"
    ], img)

    # Slide 57
    log_status("Adapter Modules - Adding adapter diagram")
    img = create_architecture_diagram("Adapter", ["Input", "Down\nProject", "Non-linear", "Up\nProject", "Residual\n+ Output"])
    add_content_slide("Adapter Modules", [
        "Small bottleneck layers inserted in network",
        "Structure:",
        ("Down-projection: d → r", 1),
        ("Non-linearity", 1),
        ("Up-projection: r → d", 1),
        ("Residual connection", 1),
        "Inserted after attention and FFN",
        "Only adapters trained",
        "Slightly more parameters than LoRA",
        "Good for vision encoders"
    ], img)

    # Slide 58
    log_status("Prompt Tuning - Adding prompt tuning diagram")
    img = create_architecture_diagram("Prompt Tuning", ["Learnable\nPrompts", "Input\nEmbeddings", "Frozen\nModel", "Task\nOutput"])
    add_content_slide("Prompt Tuning", [
        "Learn soft prompts: Continuous embeddings",
        "Prepend learnable tokens to input",
        "Visual prompt tuning:",
        ("Learn perturbations to input image", 1),
        ("Or learn visual prefix tokens", 1),
        "Text prompt tuning:",
        ("Learn continuous prompt embeddings", 1),
        "Extremely parameter-efficient",
        "Challenges: More careful tuning needed",
        "Good for: Few-shot learning"
    ], img)

    # Slide 59
    log_status("Instruction Tuning - Adding instruction format")
    img = create_simple_diagram("Instruction Format", ["Instruction", "Input\n(Image)", "Model", "Output\n(Answer)"])
    add_content_slide("Instruction Tuning", [
        "Fine-tune on instruction-following data",
        "Format: (Instruction, Input, Output)",
        "Example:",
        ("Instruction: 'Describe the image'", 1),
        ("Input: [Image]", 1),
        ("Output: 'A photo of...'", 1),
        "Benefits:",
        ("Better instruction following", 1),
        ("More natural interactions", 1),
        ("Improved zero-shot generalization", 1),
        "Used in: LLaVA, InstructBLIP, Otter"
    ], img)

    # Slide 60
    log_status("Task-Specific Fine-tuning - Adding task examples")
    img = create_simple_diagram("Tasks", ["VQA", "Captioning", "Medical\nImaging", "OCR"])
    add_content_slide("Task-Specific Fine-tuning", [
        "Adapt to specific downstream tasks",
        "Examples:",
        ("VQA: Fine-tune on VQAv2", 1),
        ("Captioning: Fine-tune on COCO", 1),
        ("Medical: Fine-tune on medical data", 1),
        "Task-specific head may be added",
        "Data: 1K-100K task examples",
        "Often use full fine-tuning or LoRA",
        "Evaluation on task benchmark"
    ], img)

    # Slide 61
    log_status("Domain Adaptation - Adding domain examples")
    img = create_simple_diagram("Domains", ["Medical", "Scientific", "E-commerce", "Remote\nSensing"])
    add_content_slide("Domain Adaptation", [
        "Adapt to specific domain:",
        ("Medical images and reports", 1),
        ("Scientific figures and papers", 1),
        ("E-commerce products", 1),
        ("Remote sensing imagery", 1),
        "Challenges:",
        ("Domain shift from pre-training", 1),
        ("Specialized terminology", 1),
        "Solutions:",
        ("Continue pre-training on domain data", 1),
        ("Task-specific fine-tuning", 1)
    ], img)

    # Slide 62
    log_status("Fine-tuning Best Practices - Adding best practices checklist")
    img = create_simple_diagram("Best Practices", ["Lower LR", "Freeze\nLayers", "Use LoRA", "Early\nStopping"])
    add_content_slide("Fine-tuning Best Practices", [
        "1. Start with appropriate pre-trained model",
        "2. Use lower learning rate than pre-training",
        "3. Freeze vision encoder initially if needed",
        "4. Use LoRA for large models",
        "5. Monitor validation performance",
        "6. Prevent overfitting:",
        ("Early stopping, regularization", 1),
        "7. Data augmentation for small datasets",
        "8. Gradient accumulation for larger batches"
    ], img)

    # Slide 63 - Section
    log_status("Section: Implementation & Applications")
    add_section_slide("Part V: Implementation & Applications")

    # Slide 64
    log_status("Implementation Frameworks - Adding framework logos/names")
    img = create_simple_diagram("Frameworks", ["HuggingFace\nTransformers", "PyTorch", "JAX/Flax", "LangChain"])
    add_content_slide("Implementation Frameworks", [
        "HuggingFace Transformers:",
        ("Pre-trained models: CLIP, BLIP, LLaVA", 1),
        ("Easy fine-tuning with Trainer API", 1),
        ("PEFT library for LoRA", 1),
        "PyTorch: Low-level flexibility",
        "JAX/Flax: XLA compilation, TPU support",
        "LangChain: Integration with applications",
        "OpenCLIP: Open-source CLIP training",
        "LLaVA codebase: Instruction tuning"
    ], img)

    # Slide 65
    log_status("Hardware Requirements - Adding GPU memory requirements")
    img = create_comparison_chart(
        ["CLIP ViT-B", "BLIP-2", "LLaVA-7B", "LLaVA-13B"],
        [0.4, 4, 14, 26],
        "GPU Memory Required (GB)"
    )
    add_content_slide("Hardware Requirements", [
        "GPU memory critical bottleneck",
        "Model size considerations:",
        ("CLIP ViT-B: ~400MB, any GPU", 1),
        ("BLIP-2: ~4GB, 16GB+ GPU", 1),
        ("LLaVA-7B: ~14GB, 24GB+ GPU", 1),
        ("LLaVA-13B: ~26GB, 40GB+ GPU", 1),
        "Memory reduction:",
        ("8-bit quantization", 1),
        ("LoRA reduces by 2-3x", 1),
        "Palmetto: A100 (40/80GB), H100"
    ], img)

    # Slide 66
    log_status("Palmetto Cluster Setup - Adding setup steps")
    img = create_simple_diagram("Palmetto Setup", ["Request\nGPU Node", "Load\nModules", "Create\nEnvironment", "Install\nPackages"])
    add_content_slide("Palmetto Cluster Setup", [
        "Clemson's HPC cluster",
        "GPU nodes: V100, A100, H100",
        "Setup steps:",
        ("1. Request GPU: qsub -I ...", 1),
        ("2. Load modules: CUDA, conda", 1),
        ("3. Create env: conda create -n multimodal", 1),
        ("4. Install: transformers, torch, peft", 1),
        "Storage: /scratch for datasets",
        "Jupyter via OnDemand",
        "See lab notebook for details"
    ], img)

    # Slide 67
    log_status("Code Example - Loading Models")
    img = create_simple_diagram("Load Model", ["Import\nLibraries", "Load\nProcessor", "Load\nModel", "Move to\nGPU"])
    add_content_slide("Loading Pre-trained Models", [
        "HuggingFace example:",
        "",
        "from transformers import CLIPModel, CLIPProcessor",
        "model = CLIPModel.from_pretrained(",
        "    'openai/clip-vit-base-patch32')",
        "processor = CLIPProcessor.from_pretrained(",
        "    'openai/clip-vit-base-patch32')",
        "",
        "model.to('cuda')",
        "",
        "See lab notebook for more examples!"
    ], img)

    # Slide 68
    log_status("Code Example - Inference")
    img = create_simple_diagram("Inference", ["Load\nImage", "Process\nInputs", "Model\nForward", "Get\nPredictions"])
    add_content_slide("Inference Example", [
        "CLIP zero-shot classification:",
        "",
        "image = Image.open('photo.jpg')",
        "text = ['a photo of a cat', 'a photo of a dog']",
        "",
        "inputs = processor(text=text, images=image,",
        "                   return_tensors='pt')",
        "outputs = model(**inputs)",
        "probs = outputs.logits_per_image.softmax(dim=1)",
        "",
        "See lab for BLIP, LLaVA examples"
    ], img)

    # Slide 69
    log_status("Code Example - Fine-tuning")
    img = create_lora_diagram()
    add_content_slide("Fine-tuning Example", [
        "Using LoRA:",
        "",
        "from peft import LoraConfig, get_peft_model",
        "",
        "lora_config = LoraConfig(",
        "    r=8, lora_alpha=16,",
        "    target_modules=['q_proj', 'v_proj'])",
        "",
        "model = get_peft_model(model, lora_config)",
        "",
        "trainer = Trainer(model, args, dataset)",
        "trainer.train()"
    ], img)

    # Slide 70
    log_status("Application - Visual Question Answering")
    img = create_simple_diagram("VQA Pipeline", ["Image +\nQuestion", "VLM", "Answer", "Applications"])
    add_content_slide("Application: Visual Question Answering", [
        "Task: Answer questions about images",
        "Input: Image + Question",
        "Output: Answer (text)",
        "Models: BLIP-2, LLaVA, InstructBLIP",
        "Use cases:",
        ("Accessibility for visually impaired", 1),
        ("Content moderation", 1),
        ("Educational tools", 1),
        ("Image search and retrieval", 1),
        "Datasets: VQAv2, GQA, OK-VQA"
    ], img)

    # Slide 71
    log_status("Application - Image Captioning")
    img = create_simple_diagram("Captioning Pipeline", ["Image", "Caption\nModel", "Generated\nCaption", "Use Cases"])
    add_content_slide("Application: Image Captioning", [
        "Task: Generate textual descriptions",
        "Input: Image",
        "Output: Caption (text)",
        "Models: BLIP, BLIP-2, GIT",
        "Use cases:",
        ("Alt-text for web accessibility", 1),
        ("Photo organization and search", 1),
        ("Social media content", 1),
        ("Medical report generation", 1),
        "Evaluation: BLEU, CIDEr, SPICE"
    ], img)

    # Slide 72
    log_status("Application - Text-to-Image Generation")
    img = create_architecture_diagram("Text-to-Image", ["Text\nPrompt", "Text\nEncoder", "Diffusion\nModel", "Generated\nImage"])
    add_content_slide("Application: Text-to-Image Generation", [
        "Task: Generate images from text",
        "Input: Text prompt",
        "Output: Image",
        "Models: DALL-E, Stable Diffusion, Imagen",
        "Architecture: Diffusion + text conditioning",
        ("Text encoder: CLIP or T5", 1),
        ("Diffusion: U-Net with cross-attention", 1),
        "Use cases:",
        ("Creative content generation", 1),
        ("Design and prototyping", 1),
        ("Data augmentation", 1)
    ], img)

    # Slide 73
    log_status("Application - Video Understanding")
    img = create_architecture_diagram("Video Understanding", ["Video\nFrames", "Spatial\nEncoder", "Temporal\nEncoder", "Predictions"])
    add_content_slide("Application: Video Understanding", [
        "Task: Understand and reason about video",
        "Input: Video (frames + audio)",
        "Output: Caption, answer, action",
        "Models: VideoMAE, TimeSformer, Video-LLaMA",
        "Challenges:",
        ("Temporal modeling", 1),
        ("Long sequences (memory)", 1),
        ("Multi-modal fusion (visual+audio+text)", 1),
        "Applications: Surveillance, sports, education"
    ], img)

    # Slide 74
    log_status("Challenges and Limitations")
    img = create_simple_diagram("Challenges", ["Hallucination", "Robustness", "Bias", "Cost"])
    add_content_slide("Challenges and Limitations", [
        "Hallucination: Generating false information",
        "Robustness: Sensitive to input variations",
        "Bias: Reflecting dataset biases",
        "Computational cost: Training and inference",
        "Data requirements: Large paired datasets",
        "Interpretability: Understanding decisions",
        "Safety: Potential for misuse",
        "Evaluation: Measuring true understanding"
    ], img)

    # Slide 75
    log_status("Future Directions")
    img = create_simple_diagram("Future", ["Scaling", "Efficiency", "Grounding", "Reasoning"])
    add_content_slide("Future Directions", [
        "Scaling: Larger models, more modalities",
        "Efficiency: Smaller models with better performance",
        "Grounding: Connecting to physical world",
        "Reasoning: Better logical and causal reasoning",
        "Embodied AI: Integration with robotics",
        "Multilinguality: Support more languages",
        "Scientific applications: Discovery and analysis",
        "Democratization: Making models accessible"
    ], img)

    # Slide 76
    log_status("Key Takeaways")
    img = create_diagram_multimodal_learning()
    add_content_slide("Key Takeaways", [
        "Multimodal learning enables richer AI",
        "Modern architectures leverage pre-trained components",
        "Contrastive learning is key for alignment",
        "Fine-tuning is practical for most applications",
        "Parameter-efficient methods make it accessible",
        "Many open-source models available",
        "Applications span VQA to text-to-image",
        "Exciting field with rapid progress"
    ], img)

    # Slide 77
    log_status("Resources")
    img = create_simple_diagram("Resources", ["Papers", "Code", "Datasets", "Tutorials"])
    add_content_slide("Resources", [
        "Papers:",
        ("CLIP, BLIP-2, LLaVA, Flamingo", 1),
        "Code:",
        ("HuggingFace Transformers, PEFT", 1),
        ("OpenCLIP, LLaVA GitHub", 1),
        "Datasets:",
        ("COCO, Visual Genome, LAION", 1),
        "Tutorials:",
        ("HuggingFace docs", 1),
        ("Papers with Code", 1),
        "Next: Lab session!"
    ], img)

    # Slide 78
    log_status("Questions & Discussion")
    img = create_diagram_multimodal_learning()
    add_title_slide("Questions & Discussion", "", img)

    print("\n" + "="*70)
    print(f"COMPLETED! Generated {slide_count} slides with images")
    print("="*70)

    return prs

if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/experiments/Multimodal_LLM_Lecture.pptx"
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    print(f"File size: {os.path.getsize(output_path) / (1024*1024):.2f} MB")
