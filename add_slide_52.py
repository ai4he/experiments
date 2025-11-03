#!/usr/bin/env python3
"""Add slide 52"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import sys

# Import diagram functions
exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

# Load existing presentation
prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-51.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

print("[Slide 52/78] Adding: CLIP Zero-Shot Classification Exercise")

# SLIDE 52: CLIP Zero-Shot Classification Exercise
img = create_simple_diagram("CLIP Exercise", ["Load\nModel", "Prepare\nImages", "Create\nPrompts", "Classify"])
add_comprehensive_slide("Lab Exercise 1: CLIP Zero-Shot Classification", [
    "Exercise Overview and Setup: In this first hands-on exercise, you'll implement zero-shot image classification using CLIP. You will load the pre-trained CLIP ViT-L/14 model from HuggingFace, prepare a set of test images, design text prompts for various categories, and perform classification without any training. The goal is to understand how CLIP's aligned embedding space enables classification on arbitrary categories through natural language descriptions. Code skeleton is provided in the notebook with TODO sections for you to complete. Estimated time: 15 minutes. Required imports: from transformers import CLIPProcessor, CLIPModel; import torch; from PIL import Image.",

    "Implementation Steps: (1) Load CLIP model and processor: model = CLIPModel.from_pretrained('openai/clip-vit-large-patch14'); processor = CLIPProcessor.from_pretrained('openai/clip-vit-large-patch14'). Move model to GPU if available. (2) Load or download test images - use provided URLs or your own images. (3) Define classification categories: categories = ['dog', 'cat', 'car', 'airplane', 'person']. (4) Create text prompts: prompts = [f'a photo of a {c}' for c in categories]. Experiment with different prompt templates. (5) Process inputs: inputs = processor(text=prompts, images=image, return_tensors='pt', padding=True). (6) Get predictions: outputs = model(**inputs); logits_per_image = outputs.logits_per_image; probs = logits_per_image.softmax(dim=1). (7) Display results: show image with predicted category and confidence scores.",

    "Prompt Engineering Experiments: Once basic classification works, experiment with prompt engineering to improve accuracy: (1) Domain-specific prompts: 'a satellite photo of {x}' for aerial imagery, 'a medical image showing {x}' for radiology. (2) Attribute-rich prompts: 'a photo of a {x}, which is a type of animal' vs. 'a photo of a {x}, which is a type of vehicle' to provide contextual disambiguation. (3) Ensemble prompts: create multiple prompt templates per category and average their embeddings: templates = ['a photo of a {}', 'a picture of a {}', 'an image showing a {}']; for each category, encode all templates and average. (4) Negative prompts: explicitly describe what the image is NOT to improve discrimination in ambiguous cases. Document which prompt strategies work best for your test images.",

    "Analysis and Discussion Questions: After completing the exercise, reflect on: (1) Which types of images does CLIP classify accurately? Which types does it struggle with? Consider factors like image quality, viewpoint, occlusion, fine-grained categories. (2) How does prompt wording affect accuracy? Compare 'dog' vs. 'a photo of a dog' vs. 'a golden retriever'. (3) What happens with out-of-distribution images (e.g., abstract art, unusual angles)? Does CLIP fail gracefully or make confident wrong predictions? (4) Limitations observed: Can CLIP count objects? Handle spatial relationships ('dog to the left of cat')? Distinguish fine-grained categories (dog breeds)? (5) Compare CLIP's zero-shot performance to what you'd expect from a supervised model trained on ImageNet. These observations prepare you for understanding model capabilities and limitations in real applications."
], img)

# Save
output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-52.pptx'
prs.save(output_file)
print(f"✓ Slide 52 completed and saved to: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
