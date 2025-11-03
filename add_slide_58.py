#!/usr/bin/env python3
"""Add slide 58"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-57.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

print("[Slide 58/78] Adding: Real-World Case Studies")

img = create_simple_diagram("Case Studies", ["Google\nLens", "GitHub\nCopilot", "Medical\nAI", "Robotics"])
add_comprehensive_slide("Real-World Case Studies", [
    "Google Lens - Visual Search at Scale: Google Lens is a production multimodal system enabling visual search and object recognition through smartphone cameras. Architecture: Mobile app captures image → image uploaded to cloud → vision model (likely evolved from Inception/EfficientNet) extracts features → multimodal retrieval matches against massive database of images with associated metadata, products, places. Capabilities: (1) Object identification - point camera at plant, animal, landmark, product to get information. (2) Text recognition and translation - OCR on signs, menus, documents with real-time translation. (3) Shopping - identify products and find purchasing options. (4) Homework help - scan math problems for solutions. Technical challenges solved: (1) Low-latency inference - sub-second response required for good UX, achieved through model optimization, edge processing, and caching. (2) Massive scale - billions of queries daily, requiring highly optimized serving infrastructure. (3) Continuous learning - models updated regularly with new data, A/B tested before deployment. (4) Multimodal integration - combining vision with text, knowledge graphs, and location data. Lessons: Production multimodal systems require end-to-end optimization from mobile capture to cloud inference to user experience.",

    "Medical Imaging AI - FDA-Approved Diagnostic Tools: Several multimodal AI systems have FDA approval for clinical use. Example: IDx-DR for diabetic retinopathy screening. System: Takes retinal fundus photographs → deep learning model (trained on 130,000 images) analyzes for signs of retinopathy → provides diagnostic decision (refer/don't refer). Unique aspects: (1) Regulatory approval - rigorous validation on diverse patient populations, >87% sensitivity and >90% specificity required. (2) High-stakes decision-making - false negatives can lead to blindness, false positives waste healthcare resources. (3) Interpretability requirements - clinicians need to understand AI reasoning, leading to attention visualization and saliency maps. (4) Bias mitigation - careful attention to demographic representation in training data to avoid disparities in care. (5) Human-in-the-loop - AI assists but doesn't replace clinicians. Implementation challenges: (1) Integration with clinical workflows and electronic health records. (2) Liability and responsibility questions. (3) Continuous monitoring for performance degradation or distribution shift as patient populations evolve. (4) Physician trust and adoption - requires transparent communication of capabilities and limitations. Lessons: Deploying AI in healthcare requires exceptional rigor in validation, transparency, fairness, and integration with existing clinical practice.",

    "Autonomous Vehicles - Embodied Multimodal AI: Self-driving cars are perhaps the most complex deployed multimodal systems, integrating vision (cameras), LiDAR (3D point clouds), radar, GPS, IMU (inertial measurement), and HD maps. Waymo/Cruise architecture: Sensor fusion combines data from all modalities → perception models detect and track objects (vehicles, pedestrians, cyclists, traffic signs, lane markings) → prediction models forecast future trajectories of dynamic objects → planning module determines vehicle actions → control system executes. Multimodal integration is critical: cameras provide semantic information and long-range detection; LiDAR provides accurate 3D geometry and works in darkness; radar detects velocity and works through fog/rain. Technical innovations: (1) Temporal modeling - tracking objects across frames, predicting trajectories seconds into the future. (2) Uncertainty quantification - model confidence affects decision-making, erring on side of caution. (3) Sim-to-real transfer - training in simulation then adapting to real world. (4) Continual learning - fleet learning where all vehicles contribute to shared learning. Open challenges: (1) Long-tail events - rare but critical scenarios (emergency vehicle, construction zone, aggressive drivers) underrepresented in training data. (2) Explainability and debugging - understanding why system made particular decision. (3) Validation - proving safety across infinite possible scenarios. Lessons: Safety-critical applications require redundancy, conservative decision-making, extensive testing, and graceful degradation.",

    "GitHub Copilot - Code Generation with Context: GitHub Copilot uses large language models (Codex, GPT-4) for code completion and generation. While primarily text-based, it exhibits multimodal reasoning by understanding code structure, documentation, and repository context. Deployment architecture: IDE plugin captures code context (current file, imported libraries, comments) → sends to API → LLM generates completion suggestions → user accepts/rejects/edits. Key features: (1) Context-aware suggestions - understands coding patterns, variable names, function signatures. (2) Multi-line generation - can write entire functions from docstrings. (3) Test generation - writes unit tests given function implementation. (4) Documentation - generates comments and docstrings. (5) Code translation - converting between programming languages. Production challenges: (1) Latency - sub-second response required for smooth IDE integration. Achieved through streaming responses and speculative generation. (2) Code correctness - generated code may have bugs, requiring user review. (3) Copyright and licensing - training on open-source code raises questions about generated code provenance. (4) Security - model might suggest vulnerable code patterns. Mitigated through security-focused fine-tuning and filtering. (5) Personalization - adapting to team coding styles while maintaining privacy. Lessons: AI coding assistants augment rather than replace developers, requiring UX that enables easy review and modification. User feedback loops essential for continuous improvement."
], img)

output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-58.pptx'
prs.save(output_file)
print(f"✓ Slide 58 completed: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
