#!/usr/bin/env python3
"""
Add slides 51-60 to existing presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import sys

# Import diagram functions
exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_ORANGE = RGBColor(246, 103, 51)
CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

slide_count = 50  # Starting from 50

def log_status(msg):
    global slide_count
    slide_count += 1
    print(f"[Slide {slide_count:2d}/78] {msg}")
    sys.stdout.flush()

# Load existing presentation
prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_PRELIMINARY_50_Slides.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

print("="*80)
print("Adding Slide 51 with comprehensive content")
print("="*80)

# SLIDE 51: Practical Lab Session Overview
log_status("Lab Session Overview - hands-on implementation")
img = create_simple_diagram("Lab Structure", ["Setup\nEnvironment", "CLIP\nDemo", "BLIP-2\nDemo", "LoRA\nFine-tuning"])
add_comprehensive_slide("Lab Session Overview", [
    "Lab Structure and Objectives: The hands-on lab session provides practical experience implementing multimodal models on Clemson's Palmetto cluster. You will work through three main components over 75 minutes: (1) Environment setup and model loading (15 minutes) - configuring your workspace, installing libraries, loading pre-trained models from HuggingFace. (2) CLIP exploration (25 minutes) - zero-shot image classification, image-text retrieval, embedding visualization. (3) BLIP-2 usage (20 minutes) - image captioning, visual question answering, comparing different LLM backends. (4) LoRA fine-tuning (15 minutes) - adapting CLIP for custom classification task using parameter-efficient methods. The lab notebook is available at Multimodal_LLM_Lab.ipynb with complete code examples, explanations, and exercises.",

    "Prerequisites and Setup: Before starting, ensure you have: (1) Access to Palmetto cluster with GPU allocation (qsub -I -l select=1:ncpus=8:ngpus=1:gpu_model=a100:mem=32gb,walltime=2:00:00). (2) Conda environment with PyTorch 2.0+, transformers 4.35+, peft, datasets, pillow, matplotlib. (3) HuggingFace cache configured to /scratch to avoid home directory quota issues (export HF_HOME=/scratch/$USER/.cache/huggingface). (4) Sample images downloaded or use provided URLs. The lab uses models that fit in 16-24GB VRAM: CLIP ViT-L/14 (~1GB), BLIP-2 with FlanT5-XL (~4GB), allowing multiple students per GPU or enabling experimentation with larger models on A100s (40-80GB).",

    "Learning Outcomes: By completing this lab, you will: (1) Gain practical experience loading and using pre-trained multimodal models from HuggingFace Hub. (2) Understand how to perform zero-shot classification and retrieval with CLIP, including prompt engineering techniques. (3) Learn to generate image captions and answer visual questions using BLIP-2 with different LLM backends (OPT, FlanT5). (4) Implement parameter-efficient fine-tuning using LoRA, training only 0.1% of parameters for task adaptation. (5) Develop skills in GPU memory optimization, batch processing, and evaluation metrics. (6) Build intuition for model capabilities and limitations through hands-on experimentation. (7) Prepare for the homework assignment which extends these techniques to full fine-tuning on custom datasets.",

    "Collaboration and Resources: This is a collaborative lab - work in pairs or small groups. Ask questions, discuss observations, and share insights. Resources: (1) Lab notebook with detailed code and markdown explanations. (2) CLIP paper (Radford et al., 2021) and BLIP-2 paper (Li et al., 2023) for reference. (3) HuggingFace documentation for transformers and peft libraries. (4) TA support for debugging and conceptual questions. (5) Optional extensions for advanced students: visualizing attention maps, comparing different CLIP model sizes, experimenting with different LoRA ranks, trying quantization (8-bit/4-bit). Save your work to /scratch/$USER/lab_session/ and checkpoint regularly. At the end, we'll have a brief discussion where groups share interesting findings or challenges encountered."
], img)

# Save updated presentation
output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-51.pptx'
prs.save(output_file)
print(f"\n{'='*80}")
print(f"✓ SLIDE 51 ADDED SUCCESSFULLY")
print(f"✓ Saved to: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
print(f"{'='*80}\n")
