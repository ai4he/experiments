#!/usr/bin/env python3
"""Add slides 61-65"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-60.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

print("[Slide 61/78] Adding: Open Challenges")
img = create_simple_diagram("Challenges", ["Hallucination", "Reasoning", "Efficiency", "Safety"])
add_comprehensive_slide("Open Challenges in Multimodal AI", [
    "Hallucination and Factual Accuracy: Multimodal models frequently generate plausible but incorrect content - describing objects not present in images, making false claims about visual content, or fabricating details. This hallucination problem stems from: (1) Training on noisy web data where image-text pairs may be weakly related. (2) Language model priors overwhelming visual evidence - models may describe common scenarios rather than actual image content. (3) Lack of explicit grounding mechanisms enforcing correspondence between generated text and visual features. Mitigation approaches: retrieval-augmented generation (retrieving relevant knowledge to ground responses), contrastive decoding (penalizing outputs inconsistent with image), attribution (citing visual evidence for claims), calibration (providing confidence scores), human feedback (RLHF to reduce hallucination). Progress is ongoing but hallucination remains a major barrier to deployment in high-stakes applications.",

    "Complex Reasoning and Compositional Understanding: Current models struggle with: (1) Multi-step logical reasoning ('If A is left of B and B is left of C, where is A relative to C?'). (2) Counting and precise quantification - answering 'How many?' reliably. (3) Fine-grained spatial reasoning - understanding complex spatial relationships and configurations. (4) Temporal reasoning in videos - understanding event sequences, causality, temporal ordering. (5) Abstract and analogical reasoning. Underlying issues: Transformer architectures may lack inductive biases for systematic reasoning, training data often lacks examples of complex reasoning (most captions are simple descriptions), and models may exploit shortcuts rather than learning true reasoning. Promising directions: chain-of-thought prompting (decomposing reasoning into steps), neuro-symbolic approaches (combining neural perception with symbolic reasoning), specialized architectures for compositional generalization, synthetic data generation for reasoning tasks.",

    "Computational Efficiency and Accessibility: State-of-the-art models require enormous computational resources, limiting accessibility. GPT-4V and Gemini require infrastructure few organizations can afford. Challenges: (1) Training costs - pre-training large multimodal models costs millions of dollars in compute. (2) Inference costs - deploying models at scale requires significant ongoing infrastructure costs. (3) Memory requirements - largest models don't fit on consumer hardware. (4) Latency - interactive applications need sub-second response times. Solutions: (1) Efficient architectures - MoE, sparse attention, optimized attention implementations. (2) Distillation - compressing large models to smaller efficient versions. (3) Quantization - 8-bit, 4-bit models with minimal quality loss. (4) PEFT methods - LoRA, adapters enabling fine-tuning on modest hardware. (5) Open-source alternatives - community-developed models democratizing access. Balancing performance and efficiency remains an active area of work.",

    "Safety, Alignment, and Responsible AI: Deploying multimodal AI responsibly requires addressing: (1) Bias and fairness - models may exhibit demographic biases, stereotypes, or perform worse on underrepresented groups. (2) Privacy - models trained on web data may memorize and reproduce private information from training images. (3) Misinformation and deepfakes - generative models can create convincing fake images, videos, facilitating deception. (4) Harmful content - models may generate violent, hateful, or explicit content despite safety filters. (5) Misuse - dual-use concerns where beneficial technologies enable malicious applications. (6) Alignment - ensuring models behave according to human values across diverse cultures and contexts. Approaches: careful data curation and filtering, bias measurement and mitigation, watermarking generated content, robust content moderation systems, red-teaming to identify failures, transparency about capabilities and limitations, stakeholder engagement in development, regulatory frameworks. No single solution suffices - responsible AI requires ongoing multifaceted efforts."
], img)

print("[Slide 62/78] Adding: Research Resources")
img = create_simple_diagram("Resources", ["Papers", "Code", "Datasets", "Community"])
add_comprehensive_slide("Research Resources and Community", [
    "Key Papers and Publications: Essential papers to read: (1) 'Attention Is All You Need' (Vaswani et al., 2017) - foundational Transformer architecture. (2) 'CLIP: Learning Transferable Visual Models From Natural Language Supervision' (Radford et al., 2021) - contrastive vision-language pre-training. (3) 'BLIP-2: Bootstrapping Language-Image Pre-training with Frozen Image Encoders and Large Language Models' (Li et al., 2023) - efficient multimodal architecture. (4) 'Flamingo: a Visual Language Model for Few-Shot Learning' (Alayrac et al., 2022) - in-context learning for vision-language. (5) 'LoRA: Low-Rank Adaptation of Large Language Models' (Hu et al., 2021) - parameter-efficient fine-tuning. Conferences: NeurIPS, ICML, ICLR (general ML), CVPR, ICCV, ECCV (computer vision), ACL, EMNLP (NLP). ArXiv.org for preprints - most papers appear as preprints before conference publication. Google Scholar and Semantic Scholar for literature search.",

    "Code Repositories and Libraries: Essential open-source tools: (1) HuggingFace Transformers - unified API for pre-trained models (github.com/huggingface/transformers). Thousands of models including CLIP, BLIP, LLaVA. Excellent documentation and tutorials. (2) HuggingFace PEFT - parameter-efficient fine-tuning (github.com/huggingface/peft). (3) PyTorch and TensorFlow - deep learning frameworks. (4) OpenCLIP - open-source CLIP implementation and training code (github.com/mlfoundations/open_clip). (5) LLaVA repository - vision-language assistant code and models. (6) LAION - large-scale datasets including LAION-400M, LAION-5B. Reproducibility: Most papers now include code releases. Papers With Code (paperswithcode.com) links papers to implementations and benchmarks. Model checkpoints on HuggingFace Hub enable using models without retraining.",

    "Datasets and Benchmarks: Important multimodal datasets: (1) MS-COCO - 123K images with captions, object detection, keypoints. (2) Flickr30K - 31K images with 5 captions each. (3) Visual Genome - detailed scene graphs with objects, attributes, relationships. (4) VQAv2 - 1.1M visual questions. (5) Conceptual Captions - 3.3M image-caption pairs. (6) LAION-400M/5B - large-scale web-scraped pairs. Evaluation benchmarks: ImageNet for zero-shot classification, Winoground for compositional understanding, MMMU for multimodal understanding, VQAv2 for question answering. Papers With Code benchmarks track state-of-the-art performance. New datasets released regularly - follow conferences and ArXiv.",

    "Community and Ongoing Learning: Engage with the community: (1) Twitter/X - many researchers share work, insights, and discuss papers. (2) Reddit - r/MachineLearning for discussions and paper summaries. (3) Discord/Slack communities - HuggingFace, LAION, and others have active communities. (4) Blogs - Distill.pub (excellent visualizations), HuggingFace blog, individual researcher blogs. (5) Video content - Yannic Kilcher (paper explanations), Two Minute Papers (research highlights). (6) Tutorials - HuggingFace course, fast.ai, Stanford CS231n (computer vision). Continuing education: Online courses (Coursera, Udacity), workshops at conferences, reading groups. Contribute to open source - fixing bugs, improving documentation, releasing models helps the community and builds your reputation. Reproducibility and openness are increasingly valued in ML research."
], img)

print("[Slide 63/78] Adding: Career Paths")
img = create_simple_diagram("Careers", ["Research", "Engineering", "Product", "Academia"])
add_comprehensive_slide("Career Paths in Multimodal AI", [
    "Research Scientist Roles: Research scientists advance the state of the art through novel algorithms, architectures, and insights. Responsibilities: (1) Reading and critiquing papers to stay current with research. (2) Formulating research questions and hypotheses. (3) Designing and running experiments to test ideas. (4) Analyzing results and drawing conclusions. (5) Writing papers and presenting at conferences. (6) Collaborating with other researchers and engineers. Work environments: (1) Industry research labs - Google Brain/DeepMind, Meta AI, OpenAI, Microsoft Research, NVIDIA Research. Typically require PhD, focus on both fundamental and applied research. (2) Startups - fast-paced, often more applied, equity compensation. (3) Academia - universities, focus on fundamental research, teaching, advising students. Requires PhD, tenure-track positions competitive. Skills needed: strong mathematical foundations (linear algebra, probability, optimization), programming (Python, PyTorch/TensorFlow), paper writing, presentation skills, creativity in problem formulation.",

    "Machine Learning Engineering: ML engineers build and deploy ML systems at scale. Responsibilities: (1) Implementing and optimizing models for production. (2) Building training pipelines - data preprocessing, distributed training, experiment tracking. (3) Deploying models - serving infrastructure, monitoring, A/B testing. (4) Improving system performance - latency optimization, resource efficiency. (5) Collaborating with researchers to transition research prototypes to production. Differs from research: focus on engineering excellence, scalability, reliability rather than novelty. Required skills: strong software engineering (algorithms, data structures, system design), deep learning frameworks, distributed systems, MLOps tools, cloud platforms (AWS, GCP, Azure). Roles at: tech companies (Google, Meta, Amazon, Microsoft), ML platforms (HuggingFace, Weights & Biases), AI startups, companies adopting ML (finance, healthcare, retail). Career growth: ML Engineer → Senior MLE → Staff/Principal Engineer → ML Platform Architect.",

    "Applied AI and Product Roles: Applied AI scientists/engineers work on product applications of ML. Responsibilities: (1) Identifying opportunities where ML adds product value. (2) Prototyping and validating ML solutions. (3) Working with product managers to define requirements and success metrics. (4) Collaborating with engineers to integrate ML into products. (5) Evaluating model performance and iterating based on user feedback. (6) Balancing ML capabilities with product constraints (latency, cost, user experience). Skills: both ML knowledge and product sense, communication skills for non-technical stakeholders, pragmatism (choosing appropriate models for constraints), business acumen. Roles at: product companies across industries. Related roles: Data Scientist (overlap with ML engineering, more focus on analysis), AI Product Manager (less technical, more strategy and requirements).",

    "Paths and Preparation: Traditional path: BS in CS/Math/Engineering → MS or PhD in ML/AI → research or engineering roles. Alternative paths: Self-taught via online courses and projects, bootcamps, transitioning from software engineering or other technical fields. Increasingly common as field grows. Building experience: (1) Personal projects showcasing skills - implement papers, build applications, contribute to open source. (2) Kaggle competitions demonstrating practical skills. (3) Research internships or RA positions during studies. (4) Publishing papers or blog posts demonstrating depth. (5) Building portfolio website with projects and write-ups. Job search: Specialized platforms (ML Jobs List, AI Jobs), company career pages, networking at conferences and online. Interviews typically include: coding (algorithms, data structures), ML knowledge (theory, practice), system design, past project discussions. Preparation: practice coding (LeetCode), review ML fundamentals, prepare to discuss projects in depth, study company's ML applications."
], img)

print("[Slide 64/78] Adding: Ethical Considerations")
img = create_simple_diagram("Ethics", ["Bias", "Privacy", "Transparency", "Accountability"])
add_comprehensive_slide("Ethical Considerations and Responsible AI", [
    "Identifying and Mitigating Bias: Multimodal AI systems can perpetuate harmful biases. Sources: (1) Data bias - training data may underrepresent certain groups or contain stereotypical associations. Web-scraped data reflects societal biases. (2) Measurement bias - evaluation metrics may not capture fairness. High average accuracy can mask poor performance on minorities. (3) Aggregation bias - single model for diverse populations may perform unequally. Manifestations: face recognition performing worse on darker skin tones, image search returning stereotypical gender associations, caption generation mentioning gender when irrelevant. Measurement: disaggregated evaluation across demographic groups, fairness metrics (demographic parity, equalized odds), qualitative analysis of outputs. Mitigation: balanced dataset curation, debiasing algorithms (adversarial debiasing, re-weighting), post-processing adjustments, diverse team perspectives in development. No perfect solution - ongoing monitoring and iteration needed.",

    "Privacy and Data Protection: Privacy risks in multimodal AI: (1) Training data memorization - models may reproduce training images or sensitive information. (2) Inference attacks - adversaries may infer whether specific data was in training set. (3) Personal information in visual content - faces, license plates, private documents. (4) Location privacy - geolocated images or metadata. Protections: (1) Data minimization - collect only necessary data. (2) De-identification - remove personal identifiers before training. (3) Differential privacy - adding noise to training process to provide mathematical privacy guarantees, at cost of performance. (4) Federated learning - training on decentralized data without centralizing. (5) Access controls and encryption for stored data. Regulations: GDPR in Europe, CCPA in California, other emerging frameworks govern data collection and use. Compliance required for deployment. Privacy-utility trade-offs often necessary.",

    "Transparency and Explainability: Users and stakeholders need to understand AI systems. Transparency dimensions: (1) Model transparency - architecture, training data, capabilities, limitations documented via model cards. (2) Inference transparency - explaining specific predictions. Attention visualization showing which image regions influenced output. Saliency maps highlighting important pixels. (3) Training transparency - disclosing datasets used, evaluation results, known biases. (4) Process transparency - how decisions about model development were made. Benefits: builds trust, enables debugging, satisfies regulatory requirements, allows informed decisions about deployment. Challenges: trade-offs between transparency and proprietary information, explanations may be misleading or incomplete, no consensus on what constitutes adequate explanation. Regulated sectors (healthcare, finance) often have explainability requirements.",

    "Accountability and Governance: Establishing responsibility for AI system behavior. Questions: (1) Who is liable when AI makes harmful decision? Developer, deployer, user? (2) How are AI systems audited? (3) What recourse exists for those harmed? (4) How should AI be regulated? Governance approaches: (1) Internal AI ethics boards reviewing projects. (2) External audits by third parties. (3) Regulatory frameworks - EU AI Act classifying systems by risk, proposed US frameworks. (4) Industry standards and best practices. (5) Algorithmic impact assessments before deployment. Roles: AI ethicists, trust & safety teams, policy specialists increasingly important. Documentation: maintaining decision logs, risk assessments, evaluation results enables accountability. Incident response: plans for handling AI failures, security breaches, bias discoveries. Interdisciplinary collaboration: engineers, ethicists, domain experts, affected communities should all inform AI development."
], img)

print("[Slide 65/78] Adding: Summary of Key Concepts")
img = create_simple_diagram("Summary", ["Theory", "Architectures", "Training", "Applications"])
add_comprehensive_slide("Summary of Key Concepts", [
    "Theoretical Foundations: We covered the mathematical and conceptual foundations of multimodal learning: (1) Multimodal representation learning - mapping heterogeneous inputs (images, text, audio) into shared embedding spaces enabling cross-modal reasoning. (2) Contrastive learning and InfoNCE loss - training models to align matching pairs while separating non-matching pairs, enabling zero-shot capabilities. (3) Attention mechanisms - self-attention for intra-modal processing, cross-attention for inter-modal interaction, enabling flexible information flow. (4) Transfer learning - leveraging pre-trained models and adapting to new tasks, reducing data and compute requirements. These foundations underpin all modern multimodal systems, from retrieval to generation to reasoning applications.",

    "Architectural Landscape: We explored major architectural paradigms: (1) Dual encoders (CLIP) - separate modality encoders projecting to shared space, efficient but limited interaction. (2) Cross-encoders (BLIP, ViLBERT) - deep fusion via cross-attention, rich reasoning but computationally expensive. (3) Frozen-backbone models (BLIP-2, Flamingo, LLaVA) - connecting pre-trained vision and language models via lightweight modules, parameter-efficient. (4) Native multimodal (Gemini) - training from scratch on multimodal data. Each architecture offers trade-offs in performance, efficiency, and flexibility. Choosing appropriate architecture depends on task requirements, data availability, and computational resources.",

    "Training and Adaptation: We learned both large-scale training and efficient fine-tuning: (1) Pre-training from scratch requires massive datasets (100M-5B pairs), distributed training (hundreds of GPUs), and weeks of compute - accessible primarily to well-resourced organizations. (2) Fine-tuning adapts pre-trained models to new domains and tasks - full fine-tuning updates all parameters, effective but expensive. (3) Parameter-efficient fine-tuning (PEFT) via LoRA, adapters, or prompt tuning updates <1% of parameters, democratizing adaptation of large models. (4) Practical considerations: data quality vs. quantity trade-offs, optimization strategies, memory management, evaluation metrics. Understanding training dynamics enables effective use of multimodal models in research and applications.",

    "Applications and Impact: Multimodal AI enables transformative applications: (1) Visual understanding - image classification, object detection, captioning, VQA allowing machines to perceive and reason about visual content. (2) Content generation - text-to-image models creating art, designs, educational materials from descriptions. (3) Multimodal search and retrieval - finding images from text queries and vice versa. (4) Embodied AI - robots understanding language instructions and visual scenes for manipulation and navigation. (5) Accessibility - tools helping visually or hearing impaired users through cross-modal translation. (6) Scientific and medical applications - analyzing complex multimodal data in healthcare, climate science, materials research. Deployment requires addressing robustness, bias, privacy, and safety. The field is rapidly evolving with new capabilities emerging regularly."
], img)

output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-65.pptx'
prs.save(output_file)
print(f"\n✓ Slides 61-65 completed: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
