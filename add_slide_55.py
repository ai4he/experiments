#!/usr/bin/env python3
"""Add slide 55"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-54.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

print("[Slide 55/78] Adding: Homework Assignment Overview")

img = create_simple_diagram("Homework", ["Exercise 1\nCLIP", "Exercise 2\nBLIP", "Exercise 3\nVQA", "Report"])
add_comprehensive_slide("Homework Assignment Overview", [
    "Assignment Structure and Objectives: The homework assignment consists of three progressive exercises that deepen your understanding of multimodal model fine-tuning. Each exercise involves training a model on a specific task, evaluating performance, analyzing results, and documenting findings. Total estimated time: 8-12 hours over one week. Deliverables: (1) Jupyter notebook with complete code, outputs, and analysis. (2) Written report (3-5 pages) discussing methodology, results, insights, and answers to discussion questions. (3) Trained model checkpoints uploaded to your scratch directory. Due date: One week from today. You may work individually or in pairs (specify collaborators). The assignment tests your ability to: implement fine-tuning pipelines, debug training issues, interpret evaluation metrics, and critically analyze model behavior.",

    "Exercise 1 - Fine-Tuning CLIP for Domain-Specific Classification: Fine-tune CLIP on Food-101 dataset (101 food categories, 1,000 images per category). Task: Adapt CLIP's vision encoder to improve classification accuracy on food images. Requirements: (1) Implement training pipeline with LoRA (r=8 or 16, your choice). (2) Train for at least 10 epochs with proper train/val split. (3) Implement data augmentation (random crop, flip, color jitter). (4) Track metrics: training/validation loss, accuracy, per-class accuracy. (5) Compare zero-shot baseline vs. fine-tuned performance. (6) Analyze: Which categories improve most? Which remain difficult? Does fine-tuning hurt zero-shot performance on non-food images? Test on ImageNet samples to check. (7) Experiment with at least two LoRA configurations (different ranks, target modules, or learning rates) and compare results. Submit learning curves, confusion matrix, and example predictions.",

    "Exercise 2 - Fine-Tuning BLIP for Image Captioning: Fine-tune BLIP on Flickr30K dataset (31K images with 5 captions each). Task: Improve caption quality through fine-tuning. Requirements: (1) Use BLIP-base model as starting point. (2) Implement caption generation training loop with cross-entropy loss. (3) Use LoRA for efficiency (apply to language model layers). (4) Train for 5 epochs with batch size appropriate for your GPU. (5) Evaluate with BLEU-4, METEOR, and CIDEr metrics using provided evaluation scripts. (6) Generate captions for 20 test images, showing before/after fine-tuning. (7) Analyze: Does fine-tuning produce more detailed captions? More accurate? Do you observe hallucination (describing things not in image)? How do quantitative metrics correlate with qualitative caption quality? (8) Ablation study: Compare training with different LoRA ranks or frozen vs. trainable vision encoder. Document computational costs (training time, memory usage).",

    "Exercise 3 - Fine-Tuning Vision-Language Model for VQA: Fine-tune a vision-language model on VQAv2 dataset subset (10K question-answer pairs). Task: Improve question-answering accuracy. Requirements: (1) Choose BLIP-2 or LLaVA as base model. (2) Implement VQA training: encode image+question, generate answer, compute loss. (3) Use LoRA on language model parameters. (4) Train for 3-5 epochs. (5) Evaluate VQA accuracy on held-out test set. (6) Analyze by question type: yes/no questions, counting, color, object recognition, spatial reasoning. (7) Error analysis: Identify failure modes. Does model fail to see objects? Misunderstand questions? Hallucinate? (8) Discussion: Compare pre-trained vs. fine-tuned performance. What types of questions benefit most from fine-tuning? How much training data is needed for good performance? What are remaining limitations? Written report should synthesize insights across all three exercises, discussing broader lessons about multimodal model adaptation, efficiency of PEFT methods, and practical considerations for real-world deployment."
], img)

output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-55.pptx'
prs.save(output_file)
print(f"✓ Slide 55 completed: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
