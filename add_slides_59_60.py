#!/usr/bin/env python3
"""Add slides 59-60"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-58.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = CLEMSON_PURPLE
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    return slide

print("[Slide 59/78] Adding: Section - Advanced Topics and Research")
add_section_slide("Part VI: Advanced Topics and Research Frontiers")

print("[Slide 60/78] Adding: Emerging Research Directions")
img = create_simple_diagram("Research", ["3D Vision", "Video\nLLMs", "Embodied\nAI", "Medical"])
add_comprehensive_slide("Emerging Research Directions", [
    "3D Vision and Spatial Understanding: Current multimodal models primarily work with 2D images, but many applications require 3D spatial understanding. Emerging directions: (1) 3D scene reconstruction from multiple views or single images - Neural Radiance Fields (NeRF) and Gaussian Splatting learn 3D representations from 2D views, enabling novel view synthesis and 3D understanding. (2) 3D object detection and segmentation - Point cloud processing with PointNet, Transformer-based architectures for 3D data. (3) Vision-language models with 3D grounding - understanding spatial relationships ('behind', 'above'), absolute positions, object sizes. (4) Robotics applications - manipulation requires 3D understanding of objects, grasping points, collision avoidance. Research challenges: Representing 3D data efficiently (voxels, point clouds, meshes, implicit functions), scaling to large scenes, integrating 3D with language for embodied AI.",

    "Long-Form Video Understanding: Most video models process short clips (seconds to minutes). Long-form video understanding (movies, lectures, surveillance footage) requires new approaches: (1) Hierarchical temporal modeling - processing video at multiple timescales (frames → shots → scenes → full video). (2) Memory-efficient architectures - processing hours of video without storing all frames. State-space models, memory networks, retrieval-augmented models. (3) Event understanding and temporal reasoning - identifying key events, understanding causal relationships, temporal ordering. (4) Video-language pre-training on narrated videos - leveraging naturally occurring narration in instructional videos, documentaries, movies. (5) Multimodal video-audio-text integration - combining visual content, speech, music, sound effects. Applications: Video search, content moderation, sports analytics, education (lecture understanding and Q&A), security (anomaly detection in long surveillance feeds).",

    "Embodied AI and Robotics: Integrating vision-language models into embodied agents that perceive and act in physical environments. Key research areas: (1) Vision-language-action models - models that take visual observations and language instructions as input, output robot actions. Examples: RT-1, RT-2 using vision Transformers with action tokens. (2) Sim-to-real transfer - training in simulation (cheap, safe, scalable) then deploying in real world. Domain randomization, domain adaptation techniques bridge the sim-real gap. (3) Interactive learning - robots learning from human feedback, demonstrations, corrections. (4) Multi-task learning - single model handling diverse manipulation tasks rather than task-specific policies. (5) Long-horizon planning - decomposing high-level goals ('clean the kitchen') into subtasks and low-level actions. Challenges: Sample efficiency (real-world data is expensive), safety (robots can cause damage), generalization to new environments and objects, real-time inference constraints.",

    "Domain-Specific Multimodal Models: Adapting general-purpose models to specialized domains with unique requirements: (1) Medical imaging - combining X-rays, CT, MRI with radiology reports, patient history, genomic data. Models like BiomedCLIP, Med-PaLM specifically trained on medical data with proper de-identification and privacy. Challenges: limited labeled data, need for interpretability and calibration, regulatory approval. (2) Scientific document understanding - parsing papers with text, equations, figures, tables. Models like Galactica, Minerva understand scientific notation, mathematical reasoning. (3) Satellite and remote sensing - combining satellite imagery with geospatial data, temporal data, text descriptions. Applications in agriculture, climate monitoring, disaster response. (4) Manufacturing and industrial - quality control, defect detection with multimodal data from cameras, sensors, process logs. (5) Legal and financial - document analysis combining text, tables, signatures, metadata. Each domain requires specialized data, evaluation metrics, and deployment considerations. Transfer learning from general models helps but domain-specific fine-tuning and data collection remain essential."
], img)

output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-60.pptx'
prs.save(output_file)
print(f"✓ Slides 59-60 completed: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
