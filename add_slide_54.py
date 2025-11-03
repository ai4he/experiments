#!/usr/bin/env python3
"""Add slide 54"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

exec(open('/home/user/experiments/generate_full_presentation.py').read().split('def create_presentation')[0])

CLEMSON_PURPLE = RGBColor(82, 45, 128)
DARK_GRAY = RGBColor(51, 51, 51)

prs = Presentation('/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-53.pptx')

def add_comprehensive_slide(title, paragraphs, img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CLEMSON_PURPLE

    if img:
        text_left, text_width = Inches(0.4), Inches(5.2)
        img_left, img_top, img_width = Inches(5.9), Inches(1.1), Inches(3.8)
    else:
        text_left, text_width = Inches(0.5), Inches(9)

    ctb = slide.shapes.add_textbox(text_left, Inches(1), text_width, Inches(6.2))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            ctf.add_paragraph()
        p = ctf.paragraphs[i]
        p.text = para_text
        p.font.size = Pt(11)
        p.space_before = Pt(9)
        p.space_after = Pt(6)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2

    if img:
        slide.shapes.add_picture(img, img_left, img_top, img_width)

    return slide

print("[Slide 54/78] Adding: LoRA Fine-Tuning Exercise")

img = create_lora_diagram()
add_comprehensive_slide("Lab Exercise 3: LoRA Fine-Tuning", [
    "Exercise Overview and Motivation: In this exercise, you'll fine-tune CLIP for a custom image classification task using LoRA (Low-Rank Adaptation). Rather than updating all 300M+ parameters in CLIP, LoRA trains only ~0.5M parameters (0.15% of total) by adding low-rank decomposition matrices to attention layers. This demonstrates parameter-efficient fine-tuning - adapting large pre-trained models to new domains with minimal computational cost and memory. You'll use the PEFT library from HuggingFace to apply LoRA, train on a small custom dataset (Food-101 subset), and evaluate performance improvements. Estimated time: 15 minutes training on A100. This exercise prepares you for the homework where you'll perform full-scale fine-tuning.",

    "Setup and Configuration: Install PEFT library: pip install peft. Load base CLIP model: from transformers import CLIPModel, CLIPProcessor; from peft import get_peft_model, LoraConfig, TaskType; model = CLIPModel.from_pretrained('openai/clip-vit-base-patch32'). Configure LoRA: config = LoraConfig(r=8, lora_alpha=16, target_modules=['q_proj', 'v_proj'], lora_dropout=0.1, bias='none', task_type=TaskType.FEATURE_EXTRACTION). The rank r=8 means each weight update is decomposed into two 8-dimensional matrices. lora_alpha=16 controls update magnitude. We apply LoRA only to query and value projections in attention layers. Wrap model: model = get_peft_model(model, config). Check trainable parameters: model.print_trainable_parameters() shows ~0.5M trainable out of 151M total.",

    "Training Loop Implementation: Prepare data: Use Food-101 subset with 10 classes, 100 images per class. Create DataLoader with batch_size=32. For each batch: (1) Process inputs: inputs = processor(text=text_labels, images=images, return_tensors='pt', padding=True). (2) Forward pass: outputs = model(**inputs); logits = outputs.logits_per_image. (3) Compute loss: loss = F.cross_entropy(logits, labels). (4) Backward and optimize: loss.backward(); optimizer.step(); optimizer.zero_grad(). Use AdamW optimizer with lr=1e-4, train for 5 epochs. Monitor training loss and validation accuracy. With LoRA, training is fast - approximately 2 minutes per epoch on A100. Without LoRA (full fine-tuning), it would take 10x longer and require 3x more memory. This efficiency enables rapid experimentation and iteration.",

    "Evaluation and Analysis: After training, evaluate on held-out test set. Compare three conditions: (1) Zero-shot CLIP (no fine-tuning) - baseline performance using generic food prompts. (2) LoRA fine-tuned CLIP - your trained model. (3) If time permits, full fine-tuning (update all parameters) for comparison. Metrics: accuracy, per-class precision/recall, confusion matrix. Expected results: zero-shot ~60-70%, LoRA fine-tuned ~85-90%, full fine-tuning ~88-92%. Key observations: LoRA achieves within 2-3% of full fine-tuning while training only 0.15% of parameters. Discussion questions: (1) Which food categories improved most with fine-tuning? Why? (2) Does LoRA preserve zero-shot capabilities on non-food images? (3) How does rank r affect performance and training speed? (4) When would you choose LoRA vs. full fine-tuning? This hands-on experience demonstrates that effective adaptation doesn't require massive computational resources - thoughtful methods like LoRA democratize fine-tuning large models."
], img)

output_file = '/home/user/experiments/Multimodal_LLM_Lecture_Slides_1-54.pptx'
prs.save(output_file)
print(f"✓ Slide 54 completed: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
